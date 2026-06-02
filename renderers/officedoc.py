"""
YancoRead — Office renderer (docx / pptx / xlsx → HTML).

Everything is converted to self-contained HTML (images inlined as data URIs)
so the frontend just drops it into the reading pane. Each format also yields a
navigable outline:
    docx  → headings
    pptx  → one entry per slide
    xlsx  → one entry per sheet
"""

import base64
import difflib
import html
import logging
import re
import zipfile
from pathlib import Path

logger = logging.getLogger('yancoread.officedoc')

# Safety caps so a pathological spreadsheet can't produce a 200MB page.
_XLSX_MAX_ROWS = 2000
_XLSX_MAX_COLS = 60


# ── DOCX ──────────────────────────────────────────────────────────────────────
# Map Word's Title/Subtitle paragraph styles to classes the editor styles and
# the save path maps back to those Word styles, so they round-trip both ways.
# Appended to mammoth's built-in heading/list defaults.
_MAMMOTH_STYLE_MAP = "\n".join([
    "p[style-name='Title'] => p.doc-title:fresh",
    "p[style-name='Subtitle'] => p.doc-subtitle:fresh",
])

# Known page sizes (portrait W×H in inches) for normalising a section's size to
# a friendly label the editor's Page Setup can show. Anything else → 'custom'.
_PAGE_SIZES = {
    'letter': (8.5, 11.0),
    'legal': (8.5, 14.0),
    'a4': (8.27, 11.69),
    'a3': (11.69, 16.54),
    'tabloid': (11.0, 17.0),
}


def _match_page_size(w_in, h_in) -> str:
    """Map (width, height) inches to a known size label, orientation-agnostic."""
    if not w_in or not h_in:
        return 'custom'
    lo, hi = sorted((w_in, h_in))
    for key, (pw, ph) in _PAGE_SIZES.items():
        if abs(lo - pw) < 0.06 and abs(hi - ph) < 0.06:
            return key
    return 'custom'


def _docx_page_setup(path: str):
    """Read section 0's page size / orientation / margins for the editor.

    Returns a dict (size, orientation, width_in, height_in, margins{}) or None
    if it can't be read — the editor falls back to its default page geometry.
    """
    try:
        from docx import Document
        from docx.enum.section import WD_ORIENT
        d = Document(path)
        if not d.sections:
            return None
        s = d.sections[0]

        def inch(v):
            try:
                return round(v.inches, 3)
            except Exception:
                return None

        w, h = inch(s.page_width), inch(s.page_height)
        landscape = (s.orientation == WD_ORIENT.LANDSCAPE) or bool(w and h and w > h)
        return {
            'size': _match_page_size(w, h),
            'orientation': 'landscape' if landscape else 'portrait',
            'width_in': w, 'height_in': h,
            'margins': {
                'top': inch(s.top_margin), 'bottom': inch(s.bottom_margin),
                'left': inch(s.left_margin), 'right': inch(s.right_margin),
            },
        }
    except Exception as e:                                       # pragma: no cover
        logger.debug("page-setup read failed for %s: %s", path, e)
        return None


def _read_hf(hf) -> dict:
    """Plain text + a page-number flag for one header/footer container.

    We only surface what the editor can faithfully round-trip: the literal text
    and whether a PAGE field is present (the live number is left to Word). Rich
    content — tables, images, tab-stops — is reported by the fidelity scan.
    """
    try:
        paras = list(hf.paragraphs)
    except Exception:
        return {'text': '', 'page_num': False}
    lines, page_num = [], False
    for p in paras:
        try:
            xml = p._p.xml
        except Exception:
            xml = ''
        if 'PAGE' in xml and ('instrText' in xml or 'fldSimple' in xml):
            page_num = True
        lines.append(p.text)
    return {'text': '\n'.join(lines).strip('\n'), 'page_num': page_num}


def _docx_headers_footers(path: str):
    """Section-0 primary header/footer for the editor, or None if unreadable."""
    try:
        from docx import Document
        d = Document(path)
        if not d.sections:
            return None
        s = d.sections[0]
        return {'header': _read_hf(s.header), 'footer': _read_hf(s.footer)}
    except Exception as e:                                       # pragma: no cover
        logger.debug("header/footer read failed for %s: %s", path, e)
        return None


# Mammoth renders a footnote reference inline as
#   <sup><a href="#footnote-N" id="footnote-ref-N">[N]</a></sup>
# and appends a trailing <ol> whose items carry id="footnote-N". We rewrite both
# into a compact, editable shape the save path can turn back into a real
# footnotes part — markers become <sup class="fn-ref">, the list a section.
_FN_REF_RE = re.compile(
    r'(?:<sup>\s*)*<a\s+href="#footnote-(\d+)"\s+id="footnote-ref-\d+"\s*>'
    r'.*?</a>(?:\s*</sup>)*', re.DOTALL)
_FN_LIST_RE = re.compile(
    r'<ol>\s*((?:<li id="footnote-\d+">.*?</li>\s*)+)</ol>', re.DOTALL)
_FN_ITEM_RE = re.compile(r'<li id="footnote-(\d+)">(.*?)</li>', re.DOTALL)
_FN_BACKLINK_RE = re.compile(r'\s*<a\s+href="#footnote-ref-\d+"\s*>.*?</a>', re.DOTALL)


def _transform_footnotes(body: str) -> str:
    """Rewrite mammoth's footnote markup into the editor's round-trip shape.

    The trailing footnote list becomes ``<section class="doc-footnotes">`` (the
    back-link arrows are dropped) and every in-text reference becomes
    ``<sup class="fn-ref" data-fn-id="N">N</sup>``. A document with no footnotes
    is returned unchanged.
    """
    if 'id="footnote-' not in body:
        return body

    def _section(m):
        items = []
        for it in _FN_ITEM_RE.finditer(m.group(1)):
            inner = _FN_BACKLINK_RE.sub('', it.group(2)).strip()
            items.append(f'<li class="fn-item" data-fn-id="{it.group(1)}">{inner}</li>')
        if not items:
            return m.group(0)
        return ('<section class="doc-footnotes" data-doc-footnotes="1">'
                '<ol class="fn-list">' + ''.join(items) + '</ol></section>')

    body = _FN_LIST_RE.sub(_section, body)
    body = _FN_REF_RE.sub(
        r'<sup class="fn-ref" data-fn-id="\1" contenteditable="false">\1</sup>', body)
    return body


def _docx_to_html(path: str) -> dict:
    import mammoth
    with open(path, 'rb') as fh:
        result = mammoth.convert_to_html(fh, style_map=_MAMMOTH_STYLE_MAP)
    body = _transform_footnotes(result.value)

    # Inject ids into headings and collect an outline.
    outline = []
    counter = {'n': 0}

    def _add_id(m):
        tag, attrs, inner = m.group(1), m.group(2), m.group(3)
        counter['n'] += 1
        anchor = f'h-{counter["n"]}'
        title = re.sub(r'<[^>]+>', '', inner).strip()
        level = int(tag[1])
        if title:
            outline.append({'title': title, 'anchor': anchor, 'level': level})
        return f'<{tag}{attrs} id="{anchor}">{inner}</{tag}>'

    body = re.sub(r'<(h[1-6])([^>]*)>(.*?)</\1>', _add_id, body, flags=re.DOTALL)
    # Page geometry and the primary header/footer travel together as one
    # "document structure" payload the editor threads back on save.
    page = _docx_page_setup(path)
    hf = _docx_headers_footers(path)
    if hf:
        page = (page or {})
        page.update(hf)
    ret = {
        'html': f'<article class="doc-page docx">{body}</article>',
        'outline': outline,
        'fidelity': detect_docx_fidelity(path),
        'page': page,
        'hasBackup': _docx_backup_path(path).is_file(),
    }
    # Tracked changes / comments → Review data + a markup body (the mammoth
    # render above stays the high-fidelity "Final" view). Only when present, so
    # ordinary documents pay nothing.
    if _docx_has_review(path):
        try:
            from docx import Document
            rdoc = Document(path)
            review = _docx_review(rdoc)
            if review['changes'] or review['comments']:
                ret['review'] = review
                ret['markupHtml'] = (
                    '<article class="doc-page docx trk-views">'
                    + _docx_markup_body(rdoc) + '</article>')
        except Exception as e:
            logger.debug("docx review build failed: %s", e)
    return ret


def _hf_is_simple(xml: str) -> bool:
    """True when a header/footer holds only plain text and at most a page-number
    field — exactly what the Tier-4b rebuild reproduces. Tables, images, legacy
    shapes, embedded objects, tab-stops (multi-column layouts) or any field
    other than PAGE/NUMPAGES make it rich, so the rebuild would drop something.
    """
    if re.search(r'<w:(tbl|drawing|pict)[ />]', xml):
        return False
    if '<v:' in xml or '<w:object' in xml:        # legacy VML / OLE object
        return False
    if '<w:tab' in xml:                           # tab stops → multi-column header
        return False
    codes = re.findall(r'<w:instrText[^>]*>(.*?)</w:instrText>', xml, re.DOTALL)
    codes += re.findall(r'w:instr="([^"]*)"', xml)
    for c in codes:
        token = (c.strip().upper().split() or [''])[0]
        if token not in ('PAGE', 'NUMPAGES'):
            return False
    return True


def _footnotes_are_simple(xml: str) -> bool:
    """True when footnotes.xml holds only plain-text notes — exactly what the
    Tier-4c rebuild reproduces (text + reference marks, across one or more
    paragraphs). Tables, images, legacy shapes, embedded objects, hyperlinks or
    fields make a note rich, so the rebuild would drop something.
    """
    if re.search(r'<w:(tbl|drawing|pict)[ />]', xml):
        return False
    if '<v:' in xml or '<w:object' in xml or '<w:hyperlink' in xml:
        return False
    if re.search(r'<w:(instrText|fldSimple|fldChar)[ >]', xml):
        return False
    return True


# Features the HTML-rebuild save path cannot reproduce. Detecting them lets the
# editor warn precisely and steer the user to "Save As" instead of a lossy
# overwrite. Plain inline images are NOT flagged — mammoth surfaces them and the
# converter re-embeds them, so they round-trip fine.
def detect_docx_fidelity(path: str) -> dict:
    """Scan a .docx for structure the editor can't rewrite.

    Returns ``{'lossy': bool, 'features': [str, ...]}``. Never raises — a scan
    failure simply reports nothing lost (the generic .bak warning still applies).
    """
    features: list[str] = []
    try:
        with zipfile.ZipFile(path) as z:
            names = set(z.namelist())

            def _read(name: str) -> str:
                try:
                    return z.read(name).decode('utf-8', 'replace')
                except Exception:
                    return ''

            # Headers / footers — a single plain-text primary header/footer now
            # round-trips (Tier 4b), so flag only what the rebuild still drops:
            # rich content, or first-page/even-page variants (>1 texty part).
            hdr = [_read(n) for n in names if re.match(r'word/header\d*\.xml$', n)]
            ftr = [_read(n) for n in names if re.match(r'word/footer\d*\.xml$', n)]
            hdr_texty = [x for x in hdr if '<w:t' in x]
            ftr_texty = [x for x in ftr if '<w:t' in x]
            if (len(hdr_texty) > 1 or len(ftr_texty) > 1
                    or any(not _hf_is_simple(x) for x in hdr_texty + ftr_texty)):
                features.append('Headers & footers')

            # Footnotes — plain-text notes now round-trip (Tier 4c); flag only
            # rich content. Endnotes aren't rebuilt, so any note text is lossy.
            # (Default parts hold only separators, so <w:t> means real notes.)
            fn_xml = _read('word/footnotes.xml')
            if '<w:t' in fn_xml and not _footnotes_are_simple(fn_xml):
                features.append('Footnotes')
            if '<w:t' in _read('word/endnotes.xml'):
                features.append('Endnotes')

            doc_xml = _read('word/document.xml')
            if 'word/comments.xml' in names:
                features.append('Comments')
            if re.search(r'<w:(ins|del)[ >]', doc_xml):
                features.append('Tracked changes')
            if doc_xml.count('<w:sectPr') > 1:
                features.append('Multiple sections')
            if re.search(r'<w:cols[^>]*w:num="(?:[2-9]|\d\d+)"', doc_xml):
                features.append('Columns')
            if '<w:fldChar' in doc_xml or '<w:instrText' in doc_xml:
                features.append('Fields / table of contents')
            if '<w:pict' in doc_xml or 'v:textbox' in doc_xml or 'wps:txbx' in doc_xml:
                features.append('Text boxes / shapes')
            if any(n.startswith('word/embeddings/') for n in names):
                features.append('Embedded objects')
    except Exception as e:                                   # pragma: no cover
        logger.debug("fidelity scan failed for %s: %s", path, e)
        return {'lossy': False, 'features': []}

    seen: list[str] = []
    for f in features:
        if f not in seen:
            seen.append(f)
    return {'lossy': bool(seen), 'features': seen}


# ── DOCX review (tracked changes + comments) ──────────────────────────────────
# mammoth renders the "Final" document (insertions kept, deletions dropped) and
# discards comments. This pair turns that loss into a feature: an lxml walk over
# document.xml extracts the changes/comments for a Review panel, and a small
# markup renderer emits the body with <ins>/<del> so the frontend can show
# Markup / Original views (the mammoth render stays the high-fidelity Final).
# python-docx 1.2.0 models TOP-LEVEL comments only — reply threads and resolved
# state are not exposed; we surface what's there and note the limitation.
_W = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'


def _ln(tag) -> str:
    return tag.split('}', 1)[-1] if isinstance(tag, str) else ''


def _docx_has_review(path: str) -> bool:
    """Cheap gate: does this .docx carry tracked changes or comments at all?"""
    try:
        with zipfile.ZipFile(path) as z:
            names = set(z.namelist())
            if 'word/comments.xml' in names:
                return True
            doc_xml = (z.read('word/document.xml').decode('utf-8', 'replace')
                       if 'word/document.xml' in names else '')
        return bool(re.search(r'<w:(ins|del)[ >]', doc_xml))
    except Exception:
        return False


def _docx_review(doc) -> dict:
    """Tracked changes (document order) + comments (with the text they anchor).
    Returns {'changes': [...], 'comments': [...], 'authors': [...]}"""
    comments = {}
    order = []
    try:
        for cm in doc.comments:
            cid = str(cm.comment_id)
            ts = getattr(cm, 'timestamp', None)
            comments[cid] = {'id': cid, 'author': cm.author or '',
                             'date': ts.isoformat() if ts else '',
                             'text': cm.text or '', 'quote': ''}
    except Exception as e:
        logger.debug("docx comments read failed: %s", e)

    changes = []
    active = []
    for elem in doc.element.body.iter():
        ln = _ln(elem.tag)
        if ln == 'commentRangeStart':
            cid = elem.get(_W + 'id')
            active.append(cid)
            if cid not in order:
                order.append(cid)
        elif ln == 'commentRangeEnd':
            cid = elem.get(_W + 'id')
            if cid in active:
                active.remove(cid)
        elif ln == 'ins':
            txt = ''.join(t.text or '' for t in elem.iter(_W + 't'))
            if txt:
                changes.append({'type': 'ins', 'author': elem.get(_W + 'author') or '',
                                'date': elem.get(_W + 'date') or '', 'text': txt})
        elif ln == 'del':
            txt = ''.join(t.text or '' for t in elem.iter(_W + 'delText'))
            if txt:
                changes.append({'type': 'del', 'author': elem.get(_W + 'author') or '',
                                'date': elem.get(_W + 'date') or '', 'text': txt})
        elif ln == 't' and active and elem.text:
            for cid in active:
                if cid in comments:
                    comments[cid]['quote'] += elem.text

    comment_list = [comments[c] for c in order if c in comments]
    for cid, c in comments.items():            # comments not anchored in the body
        if cid not in order:
            comment_list.append(c)

    authors = []
    for item in changes + comment_list:
        a = item.get('author')
        if a and a not in authors:
            authors.append(a)
    return {'changes': changes, 'comments': comment_list, 'authors': authors}


def _onoff(rpr, tag) -> bool:
    """A w:b / w:i / w:u toggle property is ON unless explicitly val='0'/'false'."""
    e = rpr.find(_W + tag)
    if e is None:
        return False
    return (e.get(_W + 'val') or 'true') not in ('0', 'false', 'off', 'none')


def _run_inner(r) -> str:
    out = []
    for child in r:
        ln = _ln(child.tag)
        if ln in ('t', 'delText'):
            out.append(html.escape(child.text or ''))
        elif ln == 'tab':
            out.append('&emsp;')
        elif ln in ('br', 'cr'):
            out.append('<br>')
        elif ln in ('drawing', 'pict', 'object'):
            out.append('<span class="trk-img">🖼</span>')
    inner = ''.join(out)
    if not inner:
        return ''
    rpr = r.find(_W + 'rPr')
    if rpr is not None:
        if _onoff(rpr, 'b'):
            inner = '<strong>' + inner + '</strong>'
        if _onoff(rpr, 'i'):
            inner = '<em>' + inner + '</em>'
        if rpr.find(_W + 'u') is not None and _onoff(rpr, 'u'):
            inner = '<u>' + inner + '</u>'
        if rpr.find(_W + 'strike') is not None and _onoff(rpr, 'strike'):
            inner = '<s>' + inner + '</s>'
        va = rpr.find(_W + 'vertAlign')
        if va is not None:
            val = va.get(_W + 'val')
            if val == 'superscript':
                inner = '<sup>' + inner + '</sup>'
            elif val == 'subscript':
                inner = '<sub>' + inner + '</sub>'
    return inner


def _author_attr(a) -> str:
    return ' data-author="%s" title="%s"' % (html.escape(a), html.escape(a)) if a else ''


def _inline_html(container) -> str:
    out = []
    for child in container:
        ln = _ln(child.tag)
        if ln == 'r':
            out.append(_run_inner(child))
        elif ln == 'ins':
            inner = _inline_html(child)
            if inner:
                out.append('<ins class="trk-ins"%s>%s</ins>' % (_author_attr(child.get(_W + 'author') or ''), inner))
        elif ln == 'del':
            inner = _inline_html(child)
            if inner:
                out.append('<del class="trk-del"%s>%s</del>' % (_author_attr(child.get(_W + 'author') or ''), inner))
        elif ln == 'hyperlink':
            inner = _inline_html(child)
            if inner:
                out.append('<a class="trk-link">' + inner + '</a>')
        elif ln in ('smartTag', 'sdt', 'sdtContent', 'fldSimple'):
            out.append(_inline_html(child))
        elif ln == 'commentRangeStart':
            out.append('<span class="cmt-anchor" data-cmt="%s">💬</span>' % html.escape(child.get(_W + 'id') or ''))
    return ''.join(out)


def _para_html(p) -> str:
    inner = _inline_html(p)
    tag = 'p'
    is_list = False
    ppr = p.find(_W + 'pPr')
    if ppr is not None:
        pstyle = ppr.find(_W + 'pStyle')
        if pstyle is not None:
            val = (pstyle.get(_W + 'val') or '').lower()
            m = re.search(r'heading(\d)', val)
            if m:
                tag = 'h' + m.group(1)
            elif val in ('title',):
                tag = 'h1'
        if ppr.find(_W + 'numPr') is not None:
            is_list = True
    if is_list:
        return '<p class="trk-li">%s</p>' % (inner or '&nbsp;')
    return '<%s>%s</%s>' % (tag, inner or '&nbsp;', tag)


def _table_html_docx(tbl) -> str:
    rows = []
    for tr in tbl.findall(_W + 'tr'):
        cells = []
        for tc in tr.findall(_W + 'tc'):
            inner = ''.join(_para_html(p) for p in tc.findall(_W + 'p'))
            cells.append('<td>%s</td>' % inner)
        rows.append('<tr>%s</tr>' % ''.join(cells))
    return '<table class="trk-table">%s</table>' % ''.join(rows)


def _docx_markup_body(doc) -> str:
    """Render the body with tracked changes visible (<ins>/<del>) plus comment
    anchors. Approximate but faithful for review: text, basic run formatting,
    headings, lists, tables and hyperlinks; images become a small placeholder
    (the Final/mammoth view keeps the real ones)."""
    parts = []
    for blk in doc.element.body:
        ln = _ln(blk.tag)
        if ln == 'p':
            parts.append(_para_html(blk))
        elif ln == 'tbl':
            parts.append(_table_html_docx(blk))
    return ''.join(parts)


# ── DOCX compare-with-backup (difflib redline) ────────────────────────────────
# An overwrite save copies the previous file to "<name>.docx.bak". Compare diffs
# the current document against that backup and renders a redline reusing the D1
# ins/del styling: paragraph-level diff, with a word-level diff inside paragraphs
# that were merely edited. Zero new deps (difflib is stdlib).
_WORD_RE = re.compile(r'\s+|\S+')


def _docx_para_texts(path: str) -> list:
    from docx import Document
    return [p.text for p in Document(path).paragraphs]


def _doc_tokens(paras: list) -> list:
    """Whole document as a word/space token stream with '\\n' paragraph breaks —
    a word-level diff over this gives a far cleaner redline than diffing whole
    paragraphs (which mis-pairs added/removed lines)."""
    toks = []
    for i, p in enumerate(paras):
        if i:
            toks.append('\n')
        toks.extend(_WORD_RE.findall(p))
    return toks


def _docx_backup_path(path: str) -> Path:
    return Path(str(path) + '.bak')          # matches app.py's <name>.docx.bak


def _word_redline(old_text: str, new_text: str):
    """Word-level redline of two (possibly multi-paragraph, '\\n'-joined) blocks.
    Returns (html_of_<p>_paragraphs, changed_paragraph_count)."""
    old_t = _doc_tokens(old_text.split('\n'))
    new_t = _doc_tokens(new_text.split('\n'))
    paras, buf = [], []
    state = {'dirty': False, 'changed': 0}

    def emit(tag, text):
        if not text:
            return
        if tag == 'equal' or not text.strip():
            buf.append(html.escape(text))
        elif tag == 'del':
            buf.append('<del class="trk-del">%s</del>' % html.escape(text)); state['dirty'] = True
        else:
            buf.append('<ins class="trk-ins">%s</ins>' % html.escape(text)); state['dirty'] = True

    def flush():
        paras.append('<p>%s</p>' % (''.join(buf) or '&nbsp;'))
        if state['dirty']:
            state['changed'] += 1
        buf.clear(); state['dirty'] = False

    def process(tokens, tag):
        seg = []
        for tok in tokens:
            if tok == '\n':
                if seg:
                    emit(tag, ''.join(seg)); seg = []
                flush()
            else:
                seg.append(tok)
        if seg:
            emit(tag, ''.join(seg))

    for tag, i1, i2, j1, j2 in difflib.SequenceMatcher(a=old_t, b=new_t, autojunk=False).get_opcodes():
        if tag == 'equal':
            process(old_t[i1:i2], 'equal')
        elif tag == 'delete':
            process(old_t[i1:i2], 'del')
        elif tag == 'insert':
            process(new_t[j1:j2], 'ins')
        else:
            process(old_t[i1:i2], 'del')
            process(new_t[j1:j2], 'ins')
    flush()
    return ''.join(paras), state['changed']


def _docx_compare(path: str) -> dict:
    """Redline of the current .docx against its last-saved .bak backup. A
    paragraph-level diff aligns whole added/removed/kept paragraphs cleanly;
    only edited (replaced) paragraphs get a word-level inline diff. Reuses the
    D1 ins/del styling."""
    bak = _docx_backup_path(path)
    if not bak.is_file():
        return {'backup': False}
    old = _docx_para_texts(str(bak))
    new = _docx_para_texts(path)
    parts = []
    changed = 0
    for tag, i1, i2, j1, j2 in difflib.SequenceMatcher(a=old, b=new, autojunk=False).get_opcodes():
        if tag == 'equal':
            for k in range(i1, i2):
                parts.append('<p>%s</p>' % (html.escape(old[k]) or '&nbsp;'))
        elif tag == 'delete':
            for k in range(i1, i2):
                if old[k].strip():
                    parts.append('<p><del class="trk-del">%s</del></p>' % html.escape(old[k])); changed += 1
                else:
                    parts.append('<p>&nbsp;</p>')
        elif tag == 'insert':
            for k in range(j1, j2):
                if new[k].strip():
                    parts.append('<p><ins class="trk-ins">%s</ins></p>' % html.escape(new[k])); changed += 1
                else:
                    parts.append('<p>&nbsp;</p>')
        else:   # replace → word-level diff on the joined block
            block_html, block_changed = _word_redline('\n'.join(old[i1:i2]), '\n'.join(new[j1:j2]))
            parts.append(block_html)
            changed += block_changed
    return {'backup': True, 'changed': changed,
            'html': '<article class="doc-page docx trk-views">%s</article>' % ''.join(parts)}


# ── DOCX export: HTML → Markdown / standalone HTML ────────────────────────────
# A small, dependency-free HTML→Markdown converter (uses lxml, already bundled).
# The source is the clean mammoth render, so the tag set is predictable:
# headings, paragraphs, b/i/u/s, links, lists (nested), blockquote, code/pre,
# hr, tables and images.
def _md_esc(text: str) -> str:
    return (text or '').replace('\\', '\\\\').replace('`', '\\`') \
                       .replace('*', '\\*').replace('[', '\\[')


def _md_inline(el) -> str:
    out = [_md_esc(el.text)] if el.text else []
    for child in el:
        out.append(_md_inline_el(child))
        if child.tail:
            out.append(_md_esc(child.tail))
    return ''.join(out)


def _md_inline_el(el) -> str:
    if not isinstance(el.tag, str):
        return ''
    tag = el.tag.lower()
    inner = _md_inline(el)
    if tag in ('strong', 'b'):
        return '**%s**' % inner if inner.strip() else inner
    if tag in ('em', 'i'):
        return '*%s*' % inner if inner.strip() else inner
    if tag in ('s', 'strike', 'del'):
        return '~~%s~~' % inner if inner.strip() else inner
    if tag == 'code':
        return '`%s`' % el.text_content()
    if tag == 'br':
        return '  \n'
    if tag == 'a':
        href = el.get('href') or ''
        return '[%s](%s)' % (inner or href, href) if href else inner
    if tag == 'img':
        alt = el.get('alt') or ''
        src = el.get('src') or ''
        return '![%s](embedded image)' % alt if src.startswith('data:') else '![%s](%s)' % (alt, src)
    return inner   # u / ins / span / sub / sup / font … → keep text


def _md_list(el, ordered: bool, depth: int) -> str:
    lines = []
    idx = 1
    for li in el:
        if not isinstance(li.tag, str) or li.tag.lower() != 'li':
            continue
        inline, nested = ([_md_esc(li.text)] if li.text else []), []
        for c in li:
            t = c.tag.lower() if isinstance(c.tag, str) else ''
            if t in ('ul', 'ol'):
                nested.append(_md_list(c, t == 'ol', depth + 1))
            else:
                inline.append(_md_inline_el(c))
            if c.tail:
                inline.append(_md_esc(c.tail))
        marker = ('%d.' % idx) if ordered else '-'
        lines.append('  ' * depth + marker + ' ' + ''.join(inline).strip())
        lines.extend(nested)
        idx += 1
    return '\n'.join(lines)


def _md_table(table) -> str:
    rows = []
    for tr in table.iter('tr'):
        cells = [_md_inline(c).strip().replace('|', '\\|').replace('\n', ' ')
                 for c in tr if isinstance(c.tag, str) and c.tag.lower() in ('td', 'th')]
        if cells:
            rows.append(cells)
    if not rows:
        return ''
    ncol = max(len(r) for r in rows)
    rows = [r + [''] * (ncol - len(r)) for r in rows]
    out = ['| ' + ' | '.join(rows[0]) + ' |', '| ' + ' | '.join(['---'] * ncol) + ' |']
    out += ['| ' + ' | '.join(r) + ' |' for r in rows[1:]]
    return '\n'.join(out)


def _md_block(el) -> str:
    if not isinstance(el.tag, str):
        return ''
    tag = el.tag.lower()
    if re.fullmatch(r'h[1-6]', tag):
        return '#' * int(tag[1]) + ' ' + _md_inline(el).strip()
    if tag in ('ul', 'ol'):
        return _md_list(el, tag == 'ol', 0)
    if tag == 'blockquote':
        inner = '\n\n'.join(b for b in (_md_block(c) for c in el if isinstance(c.tag, str)) if b.strip()) \
            or _md_inline(el).strip()
        return '\n'.join('> ' + ln for ln in inner.split('\n'))
    if tag == 'pre':
        return '```\n' + el.text_content().rstrip('\n') + '\n```'
    if tag == 'hr':
        return '---'
    if tag == 'table':
        return _md_table(el)
    if tag in ('div', 'article', 'section'):
        return '\n\n'.join(b for b in (_md_block(c) for c in el if isinstance(c.tag, str)) if b.strip())
    if tag == 'p':                  # Word Title/Subtitle styles → headings
        cls = el.get('class') or ''
        inner = _md_inline(el).strip()
        if 'doc-title' in cls:
            return '# ' + inner if inner else ''
        if 'doc-subtitle' in cls:
            return '## ' + inner if inner else ''
        return inner
    return _md_inline(el).strip()   # anything else → inline content


def html_to_markdown(html_str: str) -> str:
    from lxml import html as lhtml
    try:
        root = lhtml.fromstring(html_str)
    except Exception:
        return ''
    kids = [c for c in root if isinstance(c.tag, str)]
    sources = kids if kids else [root]
    blocks = [b for b in (_md_block(c) for c in sources) if b and b.strip()]
    return '\n\n'.join(blocks) + '\n'


def accept_reject_changes(src: str, dest: str, mode: str = 'accept') -> dict:
    """Accept or reject ALL tracked changes, writing a new .docx (the original
    is never touched). Handles the text revisions — insertions, deletions and
    moves; accepting also clears formatting-change records. python-docx/lxml,
    no new deps.

    accept: keep insertions, drop deletions.   reject: drop insertions, keep
    (restore) deletions. Comments are left intact.
    """
    from lxml import etree
    Wq = _W   # '{...wordprocessingml...}'

    with zipfile.ZipFile(src) as z:
        doc_xml = z.read('word/document.xml')
    root = etree.fromstring(doc_xml)

    def unwrap(el):
        parent = el.getparent()
        if parent is None:
            return
        idx = parent.index(el)
        for child in list(el):
            el.remove(child)
            parent.insert(idx, child); idx += 1
        parent.remove(el)

    def remove(el):
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)

    keep = ('ins', 'moveTo') if mode == 'accept' else ('del', 'moveFrom')
    drop = ('del', 'moveFrom') if mode == 'accept' else ('ins', 'moveTo')

    count = 0
    for tag in drop:                                   # remove rejected/deleted runs
        for el in root.findall('.//' + Wq + tag):
            if el.getparent() is not None:
                remove(el); count += 1
    for tag in keep:                                   # unwrap kept runs into the text
        for el in root.findall('.//' + Wq + tag):
            if el.getparent() is None:
                continue
            if mode == 'reject' and tag in ('del', 'moveFrom'):
                for dt in el.findall('.//' + Wq + 'delText'):
                    dt.tag = Wq + 't'                  # restore deleted text as live text
            unwrap(el); count += 1
    if mode == 'accept':                               # accept formatting-change records
        for tag in ('rPrChange', 'pPrChange', 'tblPrChange', 'trPrChange',
                    'tcPrChange', 'sectPrChange'):
            for el in root.findall('.//' + Wq + tag):
                remove(el)

    new_xml = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=True)
    with zipfile.ZipFile(src) as zin, \
            zipfile.ZipFile(dest, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = new_xml if item.filename == 'word/document.xml' else zin.read(item.filename)
            zout.writestr(item, data)
    return {'ok': True, 'mode': mode, 'changed': count}


def html_to_standalone(body_html: str, title: str = '') -> str:
    return (
        '<!doctype html>\n<html lang="en"><head><meta charset="utf-8">\n'
        '<meta name="viewport" content="width=device-width, initial-scale=1">\n'
        '<title>%s</title>\n<style>\n'
        'body{max-width:760px;margin:2rem auto;padding:0 1rem;color:#1a2230;'
        'font:16px/1.6 -apple-system,Segoe UI,system-ui,sans-serif;background:#fff}\n'
        'h1,h2,h3{line-height:1.25}img{max-width:100%%;height:auto}\n'
        'table{border-collapse:collapse;margin:1em 0}td,th{border:1px solid #ccd;padding:5px 9px}\n'
        'blockquote{border-left:3px solid #ccd;margin:1em 0;padding:.2em 1em;color:#555}\n'
        'code{background:#f0f2f5;padding:.1em .35em;border-radius:3px}\n'
        'pre{background:#f6f8fa;padding:12px 14px;border-radius:8px;overflow:auto}\n'
        '</style>\n</head>\n<body>\n%s\n</body></html>\n'
        % (html.escape(title or 'Document'), body_html))


# ── PPTX ──────────────────────────────────────────────────────────────────────
# Tier-0 fidelity: render each slide as absolutely-positioned shapes inside a
# native-pixel surface (the slide is `slide_size` px; the frontend scales the
# whole surface to fit). Geometry comes straight from each shape's EMU box
# (EMU → px at 96 dpi = EMU / 9525). Text boxes, pictures, tables, solid fills,
# rotation and (recursively) grouped shapes are placed; theme/scheme colors and
# gradients are skipped rather than guessed. This is a faithful-but-approximate
# native render — not a pixel-perfect rasterization (that's the optional
# LibreOffice path, P4).
_EMU_PER_PX = 9525.0
# Image content types a browser can render inline. EMF/WMF/TIFF metafiles are
# skipped (they'd show as a broken-image icon); the LibreOffice path (P4) can
# rasterize them later.
_WEB_IMAGE_TYPES = frozenset({
    'image/png', 'image/jpeg', 'image/jpg', 'image/gif',
    'image/bmp', 'image/webp', 'image/svg+xml', 'image/x-png',
})


def _pt_px(pt) -> float:
    return round(pt * 96.0 / 72.0, 1)


def _fill_color(fill) -> str:
    """Solid RGB fill → '#rrggbb'; None for inherited/theme/gradient fills."""
    try:
        from pptx.enum.dml import MSO_FILL_TYPE, MSO_COLOR_TYPE
        if fill.type == MSO_FILL_TYPE.SOLID:
            c = fill.fore_color
            if c.type == MSO_COLOR_TYPE.RGB:
                rgb = c.rgb
                return '#%02x%02x%02x' % (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return ''


def _font_color(font) -> str:
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE
        c = font.color
        if c is not None and c.type == MSO_COLOR_TYPE.RGB:
            rgb = c.rgb
            return '#%02x%02x%02x' % (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return ''


def _anchor_justify(tf) -> str:
    try:
        from pptx.enum.text import MSO_ANCHOR
        a = tf.vertical_anchor
        if a == MSO_ANCHOR.MIDDLE:
            return 'center'
        if a == MSO_ANCHOR.BOTTOM:
            return 'flex-end'
    except Exception:
        pass
    return 'flex-start'


def _para_align(para) -> str:
    try:
        from pptx.enum.text import PP_ALIGN
        return {PP_ALIGN.CENTER: 'center', PP_ALIGN.RIGHT: 'right',
                PP_ALIGN.JUSTIFY: 'justify'}.get(para.alignment, 'left')
    except Exception:
        return 'left'


def _run_html(run) -> str:
    txt = html.escape(run.text)
    if not txt:
        return ''
    f = run.font
    styles = []
    try:
        if f.size:
            styles.append('font-size:%spx' % _pt_px(f.size.pt))
    except Exception:
        pass
    if f.bold:
        styles.append('font-weight:700')
    if f.italic:
        styles.append('font-style:italic')
    if f.underline:
        styles.append('text-decoration:underline')
    col = _font_color(f)
    if col:
        styles.append('color:' + col)
    try:
        if f.name:
            styles.append("font-family:'%s'" % f.name.replace("'", ''))
    except Exception:
        pass
    return '<span style="%s">%s</span>' % (';'.join(styles), txt) if styles else txt


def _text_html(tf) -> str:
    out = []
    for para in tf.paragraphs:
        inner = ''.join(_run_html(r) for r in para.runs)
        if not inner and para.text.strip():
            inner = html.escape(para.text)
        style = 'margin:0;text-align:%s' % _para_align(para)
        lvl = para.level or 0
        if lvl:
            style += ';padding-left:%dpx' % (lvl * 24)
        out.append('<p style="%s">%s</p>' % (style, inner or '&nbsp;'))
    return ''.join(out)


def _table_html(table) -> str:
    rows = []
    for row in table.rows:
        cells = ''.join('<td>%s</td>' % html.escape(c.text) for c in row.cells)
        rows.append('<tr>%s</tr>' % cells)
    return '<table class="sl-table">%s</table>' % ''.join(rows)


def _group_transform(group, T):
    """Compose T (child-space EMU → slide EMU, as (ax,bx,ay,by)) with a group's
    own xfrm so the group's children map onto the slide."""
    ax, bx, ay, by = T
    try:
        from pptx.oxml.ns import qn
        xfrm = group._element.find(qn('p:grpSpPr')).find(qn('a:xfrm'))
        off, ext = xfrm.find(qn('a:off')), xfrm.find(qn('a:ext'))
        ch_off, ch_ext = xfrm.find(qn('a:chOff')), xfrm.find(qn('a:chExt'))
        ox, oy = int(off.get('x')), int(off.get('y'))
        ecx, ecy = int(ext.get('cx')), int(ext.get('cy'))
        cox, coy = int(ch_off.get('x')), int(ch_off.get('y'))
        ccx, ccy = int(ch_ext.get('cx')), int(ch_ext.get('cy'))
        kx = (ecx / ccx) if ccx else 1.0
        ky = (ecy / ccy) if ccy else 1.0
        return (ax + bx * (ox - cox * kx), bx * kx,
                ay + by * (oy - coy * ky), by * ky)
    except Exception:
        return T   # graceful: place children in the group's own space


def _emit_shape(shape, parts, T):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    st = shape.shape_type

    if st == MSO_SHAPE_TYPE.GROUP:
        T2 = _group_transform(shape, T)
        for child in shape.shapes:
            try:
                _emit_shape(child, parts, T2)
            except Exception as e:
                logger.debug("pptx group child skipped: %s", e)
        return

    ax, bx, ay, by = T
    left = round((ax + bx * (shape.left or 0)) / _EMU_PER_PX, 1)
    top = round((ay + by * (shape.top or 0)) / _EMU_PER_PX, 1)
    w = round((bx * (shape.width or 0)) / _EMU_PER_PX, 1)
    h = round((by * (shape.height or 0)) / _EMU_PER_PX, 1)
    box = 'left:%spx;top:%spx;width:%spx;height:%spx' % (left, top, w, h)
    rot = getattr(shape, 'rotation', 0) or 0
    if rot:
        box += ';transform:rotate(%sdeg)' % round(rot, 2)

    cls, inner, extra = 'sl-shape', '', ''

    if st == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            # Browsers can't render EMF/WMF/TIFF metafiles — skip rather than
            # emit a broken-image icon. (LibreOffice path, P4, rasterizes these.)
            if img.content_type not in _WEB_IMAGE_TYPES:
                logger.debug("pptx skipping non-web image %s", img.content_type)
                return
            b64 = base64.b64encode(img.blob).decode('ascii')
            inner = ('<img alt="" src="data:%s;base64,%s">'
                     % (img.content_type, b64))
            cls += ' sl-pic'
        except Exception as e:
            logger.debug("pptx image extract failed: %s", e)
            return
    elif getattr(shape, 'has_table', False):
        inner = _table_html(shape.table)
        cls += ' sl-tablebox'
    else:
        bg = ''
        try:
            bg = _fill_color(shape.fill)
        except Exception:
            bg = ''
        if bg:
            extra += ';background:%s' % bg
        has_text = False
        try:
            has_text = shape.has_text_frame and shape.text_frame.text.strip() != ''
        except Exception:
            has_text = False
        if has_text:
            tf = shape.text_frame
            inner = ('<div class="sl-text" style="justify-content:%s">%s</div>'
                     % (_anchor_justify(tf), _text_html(tf)))
        elif not bg:
            return   # nothing visible (empty placeholder, connector, etc.)

    parts.append('<div class="%s" style="%s%s">%s</div>' % (cls, box, extra, inner))


def _pptx_to_html(path: str) -> dict:
    from pptx import Presentation

    prs = Presentation(path)
    slides_html = []
    outline = []
    notes = []

    # Slide stage geometry: EMU → CSS px at 96 dpi. python-pptx leaves these
    # None on rare malformed decks; fall back to the classic 4:3 (10"×7.5").
    emu_w = prs.slide_width or 9144000
    emu_h = prs.slide_height or 6858000
    slide_size = {'width': round(emu_w / 9525), 'height': round(emu_h / 9525)}

    identity = (0.0, 1.0, 0.0, 1.0)
    for i, slide in enumerate(prs.slides, start=1):
        anchor = f'slide-{i}'
        title_text = ''
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text.strip()
        except Exception:
            pass
        outline.append({'title': title_text or f'Slide {i}', 'anchor': anchor, 'level': 1})

        # Speaker notes (guard with has_notes_slide so we never create one).
        note_text = ''
        try:
            if slide.has_notes_slide:
                note_text = (slide.notes_slide.notes_text_frame.text or '').strip()
        except Exception:
            note_text = ''
        notes.append(note_text)

        parts = []
        bg = ''
        try:
            bg = _fill_color(slide.background.fill)
        except Exception:
            bg = ''
        if bg:
            # Pick a readable default for un-themed text: PowerPoint resolves
            # run colors from the theme (which we skip), so on a dark slide
            # background uncolored text would default to dark and vanish. Tag
            # dark backgrounds so the surface flips its default text to light.
            cls = 'sl-bg'
            try:
                rr, gg, bb = int(bg[1:3], 16), int(bg[3:5], 16), int(bg[5:7], 16)
                if (0.299 * rr + 0.587 * gg + 0.114 * bb) < 128:
                    cls += ' sl-bg-dark'
            except Exception:
                pass
            parts.append('<div class="%s" style="background:%s"></div>' % (cls, bg))
        for shape in slide.shapes:
            try:
                _emit_shape(shape, parts, identity)
            except Exception as e:
                logger.debug("pptx shape skipped: %s", e)

        slides_html.append(
            f'<section class="slide" id="{anchor}">' + ''.join(parts) + '</section>')

    try:
        from renderers import libreoffice
        hifi = libreoffice.available()
    except Exception:
        hifi = False

    body = '\n'.join(slides_html)
    return {'html': f'<article class="doc-page pptx">{body}</article>',
            'outline': outline,
            'slide_size': slide_size,
            'slide_count': len(slides_html),
            'notes': notes,
            'hifi_available': hifi}


# ── XLSX ──────────────────────────────────────────────────────────────────────
# Structured render: each sheet becomes JSON (cells, merges, freeze, sizes) and
# the frontend draws a sticky-header/column grid with sheet tabs. We open the
# workbook WITHOUT read_only=True on purpose — merged_cells.ranges and
# freeze_panes are not exposed on a ReadOnlyWorksheet. A non-read-only read of a
# capped 2000×60 window is still ~1s, so the cap stays the safeguard.
def _cell_value(v, epoch):
    """Coerce an openpyxl value into a JSON primitive. Dates/times become Excel
    serial numbers so the frontend can format them with their number-format code
    via SSF (ssf.js). Numbers/bools/strings pass through."""
    import datetime as _dt
    if v is None or isinstance(v, (bool, int, float, str)):
        return v
    if isinstance(v, _dt.datetime) or isinstance(v, _dt.date):
        try:
            from openpyxl.utils.datetime import to_excel
            return to_excel(v, epoch)
        except Exception:
            return str(v)
    if isinstance(v, _dt.time):
        return (v.hour * 3600 + v.minute * 60 + v.second) / 86400.0
    return str(v)


def _argb_hex(color):
    """An explicit RGB cell color → '#rrggbb'. Theme / indexed / auto colors are
    skipped (we can't resolve the workbook theme), as are fully-transparent
    (alpha 00) fills."""
    try:
        if color is None or color.type != 'rgb':
            return ''
        rgb = color.rgb
        if not isinstance(rgb, str) or len(rgb) < 6:
            return ''
        if len(rgb) == 8 and rgb[:2] == '00':   # alpha 0 → no paint
            return ''
        return '#' + rgb[-6:].lower()
    except Exception:
        return ''


def _cell_style(cell):
    """Compact, JSON-small style flags for a cell: bold/italic/underline,
    horizontal align, font color, solid fill, and which borders are present."""
    s = {}
    f = cell.font
    if f is not None:
        if f.bold:
            s['b'] = 1
        if f.italic:
            s['i'] = 1
        if f.underline:
            s['u'] = 1
        fc = _argb_hex(f.color)
        if fc:
            s['fc'] = fc
    fill = cell.fill
    if fill is not None and getattr(fill, 'patternType', None) == 'solid':
        bg = _argb_hex(fill.fgColor)
        if bg:
            s['bg'] = bg
    al = cell.alignment
    if al is not None and al.horizontal in ('left', 'center', 'right'):
        s['a'] = al.horizontal[0]
    b = cell.border
    if b is not None:
        bd = ''
        for side, ch in (('top', 't'), ('bottom', 'b'), ('left', 'l'), ('right', 'r')):
            sd = getattr(b, side, None)
            if sd is not None and sd.style:
                bd += ch
        if bd:
            s['bd'] = bd
    return s


def _xlsx_to_html(path: str) -> dict:
    from openpyxl import load_workbook
    from openpyxl.utils import column_index_from_string

    wb = load_workbook(path, data_only=True)
    epoch = getattr(wb, 'epoch', None)
    # Second pass for formula strings (data_only loses them). We drive cell
    # presence from this pass so a formula cell with NO cached value still shows
    # (its formula), rather than vanishing. If it fails, we degrade to value-only.
    try:
        wbf = load_workbook(path, data_only=False)
    except Exception:
        wbf = None
    sheets = []
    outline = []

    for ws in wb.worksheets:
        anchor = 'sheet-' + re.sub(r'\W+', '-', ws.title)
        outline.append({'title': ws.title, 'anchor': anchor, 'level': 1})

        full_rows = ws.max_row or 1
        full_cols = ws.max_column or 1
        rows = min(full_rows, _XLSX_MAX_ROWS)
        cols = min(full_cols, _XLSX_MAX_COLS)
        truncated = full_rows > _XLSX_MAX_ROWS or full_cols > _XLSX_MAX_COLS

        wsf = wbf[ws.title] if (wbf is not None and ws.title in wbf.sheetnames) else None
        cells = []
        src = wsf or ws   # iterate the formula sheet when available
        for row in src.iter_rows(min_row=1, max_row=rows, min_col=1, max_col=cols):
            for scell in row:
                r, c = scell.row, scell.column
                dcell = ws.cell(row=r, column=c)        # cached value + style
                is_formula = wsf is not None and scell.data_type == 'f'
                if dcell.value is None and not is_formula:
                    continue
                entry = {'r': r, 'c': c,
                         'v': _cell_value(dcell.value, epoch),
                         'z': dcell.number_format}
                style = _cell_style(dcell)
                if style:
                    entry['s'] = style
                if is_formula:
                    entry['f'] = str(scell.value)        # e.g. '=SUM(A1:A2)'
                cells.append(entry)

        merged = []
        for rng in ws.merged_cells.ranges:
            if rng.min_row <= rows and rng.min_col <= cols:
                merged.append({'r': rng.min_row, 'c': rng.min_col,
                               'rs': min(rng.max_row, rows) - rng.min_row + 1,
                               'cs': min(rng.max_col, cols) - rng.min_col + 1})

        # Column widths (Excel "char" units → px) and row heights (pt → px),
        # only for dimensions the file actually defines (don't fabricate keys).
        col_w = {}
        for letter, dim in ws.column_dimensions.items():
            if dim.width:
                idx = column_index_from_string(letter)
                if idx <= cols:
                    col_w[idx] = round(dim.width * 7 + 5)
        row_h = {}
        for idx, dim in ws.row_dimensions.items():
            if dim.height and idx <= rows:
                row_h[idx] = round(dim.height * 96 / 72)

        sheets.append({
            'name': ws.title, 'anchor': anchor,
            'rows': rows, 'cols': cols,
            'cells': cells, 'merged': merged,
            'freeze': ws.freeze_panes,         # e.g. 'B2' or None
            'colWidths': col_w, 'rowHeights': row_h,
            'truncated': truncated,
        })

    wb.close()
    if wbf is not None:
        wbf.close()
    return {'sheets': sheets, 'outline': outline}


# ── dispatch ────────────────────────────────────────────────────────────────
def to_html(path: str) -> dict:
    """Render an office document to {html, outline}."""
    ext = Path(path).suffix.lower()
    if ext == '.docx':
        return _docx_to_html(path)
    if ext == '.pptx':
        return _pptx_to_html(path)
    if ext == '.xlsx':
        return _xlsx_to_html(path)
    raise ValueError(f'Unsupported office format: {ext}')


# ── HTML → DOCX (editor write-back) ───────────────────────────────────────────
# A focused, native converter that maps the rich-text editor's HTML straight to
# a .docx. We own both sides (the contenteditable output and this parser), so we
# can guarantee fidelity for exactly the formatting the editor exposes — fonts,
# sizes, colours, highlight, bold/italic/underline/strike, alignment, line
# spacing, headings, lists, links, tables and inline images — with no extra
# dependency (python-docx + lxml are already present via the readers above).

_BLOCK_TAGS = {
    'p', 'div', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'blockquote', 'pre',
    'section', 'article', 'header', 'figure', 'ul', 'ol',
}
_HEADING_STYLE = {
    'h1': 'Heading 1', 'h2': 'Heading 2', 'h3': 'Heading 3',
    'h4': 'Heading 4', 'h5': 'Heading 5', 'h6': 'Heading 6',
}


def _style_map(el) -> dict:
    raw = el.get('style') or ''
    out = {}
    for part in raw.split(';'):
        if ':' in part:
            k, v = part.split(':', 1)
            out[k.strip().lower()] = v.strip()
    return out


def _parse_color(v):
    from docx.shared import RGBColor
    if not v:
        return None
    v = v.strip().lower()
    try:
        if v.startswith('#'):
            h = v[1:]
            if len(h) == 3:
                h = ''.join(c * 2 for c in h)
            if len(h) >= 6:
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        m = re.match(r'rgba?\(([^)]+)\)', v)
        if m:
            nums = [p.strip() for p in m.group(1).split(',')]
            r, g, b = (int(round(float(nums[i]))) for i in range(3))
            return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))
        named = {
            'black': '000000', 'white': 'ffffff', 'red': 'ff0000', 'green': '008000',
            'blue': '0000ff', 'yellow': 'ffff00', 'gray': '808080', 'grey': '808080',
            'orange': 'ffa500', 'purple': '800080',
        }
        if v in named:
            h = named[v]
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return None
    return None


def _parse_pt(v):
    if not v:
        return None
    v = v.strip().lower()
    try:
        if v.endswith('pt'):
            return float(v[:-2])
        if v.endswith('px'):
            return round(float(v[:-2]) * 0.75, 1)
        if v.endswith('em'):
            return round(float(v[:-2]) * 12.0, 1)
        if v.endswith('%'):
            return round(float(v[:-1]) / 100.0 * 12.0, 1)
        return float(v)
    except ValueError:
        return None


def _ctx_from(el, ctx: dict) -> dict:
    """Derive an inline-formatting context from an element's tag + inline CSS."""
    c = dict(ctx)
    tag = el.tag.lower() if isinstance(el.tag, str) else ''
    if tag in ('b', 'strong'):
        c['bold'] = True
    if tag in ('i', 'em'):
        c['italic'] = True
    if tag == 'u':
        c['underline'] = True
    if tag in ('s', 'strike', 'del'):
        c['strike'] = True
    if tag == 'sup':
        c['sup'] = True
    if tag == 'sub':
        c['sub'] = True
    if tag == 'font':
        if el.get('face'):
            c['font'] = el.get('face')
        if el.get('color'):
            col = _parse_color(el.get('color'))
            if col is not None:
                c['color'] = col

    sd = _style_map(el)
    fw = sd.get('font-weight', '')
    if fw == 'bold' or (fw.isdigit() and int(fw) >= 600):
        c['bold'] = True
    elif fw in ('normal', '400'):
        c['bold'] = False
    fs = sd.get('font-style')
    if fs == 'italic':
        c['italic'] = True
    elif fs == 'normal':
        c['italic'] = False
    deco = (sd.get('text-decoration', '') + ' ' + sd.get('text-decoration-line', ''))
    if 'underline' in deco:
        c['underline'] = True
    if 'line-through' in deco:
        c['strike'] = True
    if sd.get('color'):
        col = _parse_color(sd['color'])
        if col is not None:
            c['color'] = col
    bg = sd.get('background-color') or sd.get('background')
    if bg:
        col = _parse_color(bg)
        if col is not None:
            c['highlight'] = col
    if sd.get('font-family'):
        c['font'] = sd['font-family'].split(',')[0].strip().strip('\'"')
    if sd.get('font-size'):
        pt = _parse_pt(sd['font-size'])
        if pt:
            c['size'] = pt
    va = sd.get('vertical-align')
    if va == 'super':
        c['sup'] = True
    elif va == 'sub':
        c['sub'] = True
    return c


def _shade_run(run, rgb):
    """Apply a background fill (text highlight) to a run via w:shd."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    rpr = run._element.get_or_add_rPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), str(rgb))
    rpr.append(shd)


def _add_run(p, text, ctx):
    from docx.shared import Pt
    if not text:
        return
    # Drop pure-whitespace segments that are just HTML pretty-printing.
    if not text.strip() and ('\n' in text or '\t' in text):
        return
    text = text.replace('\n', ' ').replace('\t', ' ').replace('\xa0', ' ')
    r = p.add_run(text)
    f = r.font
    if ctx.get('bold'):
        r.bold = True
    if ctx.get('italic'):
        r.italic = True
    if ctx.get('underline'):
        r.underline = True
    if ctx.get('strike'):
        f.strike = True
    if ctx.get('font'):
        f.name = ctx['font']
    if ctx.get('size'):
        f.size = Pt(ctx['size'])
    if ctx.get('color') is not None:
        f.color.rgb = ctx['color']
    if ctx.get('sup'):
        f.superscript = True
    if ctx.get('sub'):
        f.subscript = True
    if ctx.get('highlight') is not None:
        _shade_run(r, ctx['highlight'])


_CONTENT_WIDTH_IN = 6.5      # usual letter page minus 1" margins


def _css_length(value):
    """Parse a CSS length (``%``, ``px``, ``cm``, ``in`` or bare px) into a docx
    Length, or None. Percent is taken as a fraction of the text column width."""
    from docx.shared import Cm, Inches
    w = (value or '').strip().lower()
    if not w:
        return None
    try:
        if w.endswith('%'):
            return Inches(_CONTENT_WIDTH_IN * float(w[:-1]) / 100.0)
        if w.endswith('px'):
            return Inches(float(w[:-2]) / 96.0)
        if w.endswith('cm'):
            return Cm(float(w[:-2]))
        if w.endswith('in'):
            return Inches(float(w[:-2]))
        return Inches(float(w) / 96.0)
    except ValueError:
        return None


def _img_width(el):
    """Explicit display width for an <img>, from CSS width or the width attr."""
    return _css_length(_style_map(el).get('width') or el.get('width') or '')


def _set_image_alt(shape, alt):
    """Set an inline image's alt text (the drawing's docPr descr/title)."""
    from docx.oxml.ns import qn
    docPr = shape._inline.find(qn('wp:docPr'))
    if docPr is not None:
        docPr.set('descr', alt)
        docPr.set('title', alt)


def _add_image(p, el):
    import base64
    import io
    from docx.shared import Inches
    src = el.get('src', '')
    if not src.startswith('data:'):
        return
    try:
        _, b64 = src.split(',', 1)
        data = base64.b64decode(b64)
        run = p.add_run()
        shape = run.add_picture(io.BytesIO(data))
        want = _img_width(el)
        max_w = Inches(6.0)
        if want is not None and shape.width:
            ratio = int(want) / shape.width
            shape.height = int(shape.height * ratio)
            shape.width = int(want)
        elif shape.width and shape.width > max_w:
            ratio = max_w / shape.width
            shape.height = int(shape.height * ratio)
            shape.width = max_w
        alt = (el.get('alt') or '').strip()
        if alt:
            _set_image_alt(shape, alt)
    except Exception as e:
        logger.debug("docx image embed failed: %s", e)


_HYPERLINK_RELTYPE = (
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink'
)


def _add_hyperlink(p, url, el, ctx):
    """Insert a real Word hyperlink (``w:hyperlink`` + external relationship)
    wrapping the link's inline content. Falls back to leaving the styled runs
    in place if the relationship can't be created."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import RGBColor

    lctx = dict(ctx)
    lctx.setdefault('color', RGBColor(0x0B, 0x66, 0xC2))
    lctx['underline'] = True

    before = {id(c) for c in p._p}
    _render_inline(el, lctx, p)
    new_runs = [c for c in p._p
                if id(c) not in before and c.tag == qn('w:r')]
    if not new_runs:
        return
    try:
        r_id = p.part.relate_to(url, _HYPERLINK_RELTYPE, is_external=True)
    except Exception as e:                       # pragma: no cover - defensive
        logger.debug("hyperlink relationship failed: %s", e)
        return  # the runs are already styled like a link — good enough
    hyperlink = OxmlElement('w:hyperlink')
    hyperlink.set(qn('r:id'), r_id)
    new_runs[0].addprevious(hyperlink)
    for r in new_runs:
        hyperlink.append(r)          # lxml append reparents the run element


def _place_child(p, child, ctx) -> bool:
    """Render one inline child into paragraph p. Returns False for block nodes."""
    from docx.shared import RGBColor
    if not isinstance(child.tag, str):
        return True  # comments / processing instructions
    tag = child.tag.lower()
    if tag == 'sup' and child.get('data-fn-final-id'):
        _add_footnote_ref(p, child.get('data-fn-final-id'))
    elif tag == 'br':
        p.add_run().add_break()
    elif tag == 'img':
        _add_image(p, child)
    elif tag in _BLOCK_TAGS or tag == 'table':
        return False
    elif tag == 'a':
        href = (child.get('href') or '').strip()
        if href and not href.lower().startswith('javascript:'):
            _add_hyperlink(p, href, child, ctx)
        else:
            cctx = _ctx_from(child, ctx)
            cctx.setdefault('color', RGBColor(0x0B, 0x66, 0xC2))
            cctx['underline'] = True
            _render_inline(child, cctx, p)
    else:
        _render_inline(child, _ctx_from(child, ctx), p)
    return True


def _render_inline(el, ctx, p):
    if el.text:
        _add_run(p, el.text, ctx)
    for child in el:
        _place_child(p, child, ctx)
        if child.tail:
            _add_run(p, child.tail, ctx)


def _apply_para_format(p, sd):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt
    ta = sd.get('text-align')
    align = {
        'left': WD_ALIGN_PARAGRAPH.LEFT, 'center': WD_ALIGN_PARAGRAPH.CENTER,
        'right': WD_ALIGN_PARAGRAPH.RIGHT, 'justify': WD_ALIGN_PARAGRAPH.JUSTIFY,
    }.get(ta)
    if align is not None:
        p.alignment = align
    lh = sd.get('line-height')
    if lh:
        try:
            p.paragraph_format.line_spacing = float(lh)
        except ValueError:
            pt = _parse_pt(lh)
            if pt:
                p.paragraph_format.line_spacing = Pt(pt)


# Upper bound on a single cell's span. A hostile/garbled colspan drives the
# table's column count (and thus rows×cols real cell objects), so without a cap
# a pasted <td colspan="99999"> hangs the save building 100k cells. Real tables
# never approach this.
_MAX_SPAN = 256


def _cell_spans(cell):
    """(colspan, rowspan) for an HTML cell, each clamped to 1.._MAX_SPAN."""
    def _i(name):
        try:
            return max(1, min(_MAX_SPAN, int(cell.get(name, '1'))))
        except (TypeError, ValueError):
            return 1
    return _i('colspan'), _i('rowspan')


def _cell_width(cell):
    """Explicit cell width as a docx Length, or None. Accepts %, px, cm, in."""
    return _css_length(_style_map(cell).get('width') or cell.get('width') or '')


def _set_repeat_header(row):
    """Mark a table row as a heading that repeats atop each page (w:tblHeader)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    trPr = row._tr.get_or_add_trPr()
    th = OxmlElement('w:tblHeader')
    th.set(qn('w:val'), 'true')
    trPr.append(th)


def _render_cell(html_cell, cell, base):
    """Render a td/th's content into a docx table cell.

    Editor and mammoth markup wrap cell text in block tags (``<p>``, headings,
    lists), so inline-only rendering would silently drop it. The cell's
    pre-existing empty paragraph is reused for the first block to avoid a blank
    leading line; ``base`` carries inherited formatting (e.g. bold for <th>)."""
    blocks = [c for c in html_cell
              if isinstance(c.tag, str) and c.tag.lower() in _BLOCK_TAGS]
    if not blocks:
        _render_inline(html_cell, _ctx_from(html_cell, base), cell.paragraphs[0])
        return

    used_first = False
    if (html_cell.text or '').strip():
        _add_run(cell.paragraphs[0], html_cell.text, _ctx_from(html_cell, base))
        used_first = True

    def _next_para():
        nonlocal used_first
        if not used_first:
            used_first = True
            return cell.paragraphs[0]
        return cell.add_paragraph()

    for child in html_cell:
        if not isinstance(child.tag, str):
            continue
        tag = child.tag.lower()
        if tag in ('ul', 'ol'):
            for li in child:
                if isinstance(li.tag, str) and li.tag.lower() == 'li':
                    _render_inline(li, _ctx_from(li, base), _next_para())
        elif tag in _BLOCK_TAGS:
            p = _next_para()
            _apply_para_format(p, _style_map(child))
            _render_inline(child, _ctx_from(child, base), p)
        if child.tail and child.tail.strip():
            _add_run(cell.paragraphs[-1], child.tail, base)

    if not used_first:                       # nothing rendered → fall back inline
        _render_inline(html_cell, _ctx_from(html_cell, base), cell.paragraphs[0])


def _emit_table(doc, el):
    parsed = []
    for tr in el.iter('tr'):
        cells = [c for c in tr
                 if isinstance(c.tag, str) and c.tag.lower() in ('td', 'th')]
        if cells:
            parsed.append(cells)
    if not parsed:
        return
    nrows = len(parsed)

    # Walk an occupancy grid so colspan/rowspan map to real cell coordinates.
    occupied = set()
    placements = []          # (row, col, colspan, rowspan, cell_el)
    ncols = 0
    for r, cells in enumerate(parsed):
        c = 0
        for cell in cells:
            while (r, c) in occupied:
                c += 1
            cs, rs = _cell_spans(cell)
            for dr in range(rs):
                for dc in range(cs):
                    occupied.add((r + dr, c + dc))
            placements.append((r, c, cs, rs, cell))
            c += cs
            ncols = max(ncols, c)
    if ncols < 1:
        return

    table = doc.add_table(rows=nrows, cols=ncols)
    try:
        table.style = 'Table Grid'
    except Exception:
        pass

    col_widths = {}          # grid-column → Length (from simple, unspanned cells)
    for (r, c, cs, rs, cell) in placements:
        target = table.cell(r, c)
        if cs > 1 or rs > 1:
            try:
                target = target.merge(table.cell(r + rs - 1, c + cs - 1))
            except Exception as e:                       # pragma: no cover
                logger.debug("cell merge failed: %s", e)
        base = {'bold': True} if cell.tag.lower() == 'th' else {}
        _render_cell(cell, target, base)
        w = _cell_width(cell)
        if w is not None and cs == 1 and c not in col_widths:
            col_widths[c] = w

    if col_widths:
        table.autofit = False
        for ci, width in col_widths.items():
            for ri in range(nrows):
                try:
                    table.cell(ri, ci).width = width
                except Exception:                        # pragma: no cover
                    pass

    # First row entirely <th> → a repeating header row.
    if parsed[0] and all(c.tag.lower() == 'th' for c in parsed[0]):
        _set_repeat_header(table.rows[0])


def _add_hr(doc):
    """Insert a horizontal rule as an empty paragraph with a bottom border."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '6')
    bottom.set(qn('w:space'), '1')
    bottom.set(qn('w:color'), 'auto')
    pBdr.append(bottom)
    pPr.append(pBdr)
    return p


def _emit_block(doc, el, list_style=None):
    tag = el.tag.lower() if isinstance(el.tag, str) else ''
    if tag in ('ul', 'ol'):
        style = 'List Number' if tag == 'ol' else 'List Bullet'
        for li in el:
            if isinstance(li.tag, str) and li.tag.lower() == 'li':
                _emit_block(doc, li, list_style=style)
        return
    if tag == 'table':
        _emit_table(doc, el)
        return
    if tag == 'hr':
        _add_hr(doc)
        return

    classes = set((el.get('class') or '').split())
    if list_style:
        para_style = list_style
    elif 'doc-title' in classes:
        para_style = 'Title'
    elif 'doc-subtitle' in classes:
        para_style = 'Subtitle'
    elif tag in _HEADING_STYLE:
        para_style = _HEADING_STYLE[tag]
    elif tag == 'blockquote':
        para_style = 'Quote'
    else:
        para_style = None

    try:
        p = doc.add_paragraph(style=para_style) if para_style else doc.add_paragraph()
    except KeyError:
        p = doc.add_paragraph()  # style missing in template
    _apply_para_format(p, _style_map(el))
    _render_inline(el, _ctx_from(el, {}), p)

    # Nested lists living inside this block (e.g. a <ul> inside an <li>).
    for child in el:
        if isinstance(child.tag, str) and child.tag.lower() in ('ul', 'ol'):
            _emit_block(doc, child)


def _apply_page_setup(doc, page) -> None:
    """Apply editor page-setup (size / orientation / margins) to section 0.

    python-docx does NOT swap width/height when you flip orientation, so we set
    the dimensions explicitly and keep ``w:orient`` consistent with them.
    """
    if not page:
        return
    from docx.shared import Inches
    from docx.enum.section import WD_ORIENT
    try:
        s = doc.sections[0]
    except (IndexError, AttributeError):
        return

    w = page.get('width_in')
    h = page.get('height_in')
    orient = (page.get('orientation') or '').lower()
    # Fall back to the named size when explicit dimensions weren't supplied.
    if (not w or not h) and page.get('size') in _PAGE_SIZES:
        w, h = _PAGE_SIZES[page['size']]
    if w and h:
        if orient == 'landscape' and w < h:
            w, h = h, w
        elif orient == 'portrait' and w > h:
            w, h = h, w
        s.orientation = WD_ORIENT.LANDSCAPE if (w > h) else WD_ORIENT.PORTRAIT
        s.page_width = Inches(w)
        s.page_height = Inches(h)

    m = page.get('margins') or {}
    for attr, key in (('top_margin', 'top'), ('bottom_margin', 'bottom'),
                      ('left_margin', 'left'), ('right_margin', 'right')):
        v = m.get(key)
        if isinstance(v, (int, float)) and 0 <= v <= 12:
            setattr(s, attr, Inches(v))


def _add_page_number_field(paragraph) -> None:
    """Append a live PAGE field (fldSimple form — widely supported) to a paragraph."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    fld = OxmlElement('w:fldSimple')
    fld.set(qn('w:instr'), ' PAGE ')
    r = OxmlElement('w:r')
    t = OxmlElement('w:t')
    t.text = '1'                                   # cached value; Word recomputes
    r.append(t)
    fld.append(r)
    paragraph._p.append(fld)


def _write_hf(hf, spec: dict) -> None:
    """Write plain text (+ an optional page-number field) into a header/footer.

    A fresh ``Document()`` starts with linked (inherited) headers/footers, so we
    only materialise a part when the editor actually has content — blank
    specs leave the section's default untouched.
    """
    if not isinstance(spec, dict):
        return
    text = (spec.get('text') or '').replace('\r', '')
    page_num = bool(spec.get('page_num'))
    if not text.strip() and not page_num:
        return

    hf.is_linked_to_previous = False               # give it its own part
    paras = hf.paragraphs
    first = paras[0]
    for extra in paras[1:]:                        # collapse to a clean slate
        extra._p.getparent().remove(extra._p)
    for r in list(first.runs):
        r._r.getparent().remove(r._r)

    lines = text.split('\n') if text else ['']
    cur = first
    for i, line in enumerate(lines):
        if i:
            cur = hf.add_paragraph()
        if line:
            cur.add_run(line.replace('\t', ' '))
    if page_num:
        if (lines[-1] if lines else '').strip():
            cur.add_run(' ')
        _add_page_number_field(cur)


def _apply_headers_footers(doc, page) -> None:
    """Apply the editor's primary header/footer to section 0 (plain text + page #)."""
    if not page:
        return
    try:
        s = doc.sections[0]
    except (IndexError, AttributeError):
        return
    for key, getter in (('header', lambda: s.header), ('footer', lambda: s.footer)):
        spec = page.get(key)
        if isinstance(spec, dict):
            try:
                _write_hf(getter(), spec)
            except Exception as e:                 # pragma: no cover - defensive
                logger.debug("%s write failed: %s", key, e)


_W_NS = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
_FN_REL_TYPE = ('http://schemas.openxmlformats.org/officeDocument/2006/'
                'relationships/footnotes')
_FN_CT_OVERRIDE = ('<Override PartName="/word/footnotes.xml" ContentType='
                   '"application/vnd.openxmlformats-officedocument.'
                   'wordprocessingml.footnotes+xml"/>')


def _cls(el) -> set:
    return set((el.get('class') or '').split())


def _extract_footnote_defs(frag) -> dict:
    """Pull footnote text out of the body fragment and strip the sections.

    Returns ``{orig_id: [line, ...]}`` (one entry per ``.fn-item``, paragraphs
    kept as separate lines); each ``.doc-footnotes`` section is removed so it
    never lands in the document body.
    """
    defs = {}
    for sec in list(frag.iter('section')):
        if 'doc-footnotes' not in _cls(sec):
            continue
        for li in sec.iter('li'):
            fid = li.get('data-fn-id')
            if 'fn-item' not in _cls(li) or not fid:
                continue
            blocks = [c for c in li if isinstance(c.tag, str)
                      and c.tag.lower() in _BLOCK_TAGS]
            lines = [' '.join(b.text_content().split()) for b in (blocks or [li])]
            defs[fid] = [ln for ln in lines if ln] or ['']
        parent = sec.getparent()
        if parent is not None:
            parent.remove(sec)
    return defs


def _order_footnote_refs(frag, defs) -> list:
    """Number footnote markers in document order; stamp each with a final id and
    return the ordered ``[(id, [line, ...]), ...]`` payload for footnotes.xml."""
    order = []
    for sup in frag.iter('sup'):
        if 'fn-ref' not in _cls(sup):
            continue
        fid = str(len(order) + 1)
        sup.set('data-fn-final-id', fid)
        order.append((fid, defs.get(sup.get('data-fn-id'), [''])))
    return order


def _add_footnote_ref(p, fid: str) -> None:
    """Append a superscript footnote-reference mark (Word auto-numbers it)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    r = OxmlElement('w:r')
    rpr = OxmlElement('w:rPr')
    va = OxmlElement('w:vertAlign')
    va.set(qn('w:val'), 'superscript')
    rpr.append(va)
    r.append(rpr)
    ref = OxmlElement('w:footnoteReference')
    ref.set(qn('w:id'), str(fid))
    r.append(ref)
    p._p.append(r)


def _build_footnotes_xml(order) -> str:
    """Serialise the ordered footnote payload into a footnotes.xml part. The two
    leading separator notes (ids -1 and 0) are what Word writes by default."""
    head = (
        '<w:footnote w:type="separator" w:id="-1"><w:p><w:r><w:separator/>'
        '</w:r></w:p></w:footnote>'
        '<w:footnote w:type="continuationSeparator" w:id="0"><w:p><w:r>'
        '<w:continuationSeparator/></w:r></w:p></w:footnote>')
    notes = []
    for fid, lines in order:
        paras = []
        for i, line in enumerate(lines or ['']):
            txt = html.escape(line, quote=False)
            if i == 0:
                run = ('<w:r><w:rPr><w:vertAlign w:val="superscript"/></w:rPr>'
                       '<w:footnoteRef/></w:r>'
                       f'<w:r><w:t xml:space="preserve"> {txt}</w:t></w:r>')
            else:
                run = f'<w:r><w:t xml:space="preserve">{txt}</w:t></w:r>'
            paras.append(f'<w:p>{run}</w:p>')
        notes.append(f'<w:footnote w:id="{fid}">{"".join(paras)}</w:footnote>')
    return ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<w:footnotes xmlns:w="{_W_NS}">{head}{"".join(notes)}</w:footnotes>')


def _inject_footnotes(path: str, order) -> None:
    """Splice a footnotes part into a saved .docx: add the part, a content-type
    override and a document relationship. python-docx can't model footnotes, so
    we rewrite the package zip directly."""
    import shutil
    with zipfile.ZipFile(path) as z:
        data = {n: z.read(n) for n in z.namelist()}

    ct = data.get('[Content_Types].xml', b'').decode('utf-8')
    if 'word/footnotes.xml' not in ct:
        ct = ct.replace('</Types>', _FN_CT_OVERRIDE + '</Types>')
        data['[Content_Types].xml'] = ct.encode('utf-8')

    rels_name = 'word/_rels/document.xml.rels'
    raw = data.get(rels_name)
    rels = raw.decode('utf-8') if raw else (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/'
        'package/2006/relationships"></Relationships>')
    if _FN_REL_TYPE not in rels:
        used = [int(x) for x in re.findall(r'Id="rId(\d+)"', rels)]
        nid = (max(used) + 1) if used else 1
        rel = (f'<Relationship Id="rId{nid}" Type="{_FN_REL_TYPE}" '
               'Target="footnotes.xml"/>')
        rels = rels.replace('</Relationships>', rel + '</Relationships>')
    data[rels_name] = rels.encode('utf-8')

    data['word/footnotes.xml'] = _build_footnotes_xml(order).encode('utf-8')

    tmp = path + '.tmpfn'
    with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as z:
        for name, blob in data.items():
            z.writestr(name, blob)
    shutil.move(tmp, path)


def html_to_docx(body_html: str, dest: str, page=None) -> None:
    """Convert the editor's HTML (inner content of the doc article) to a .docx.

    ``page`` is the optional document-structure dict from the editor (size /
    orientation / margins, plus the primary header & footer); when present it is
    applied to the document's section. Footnotes travel inside ``body_html`` (a
    ``.doc-footnotes`` section plus ``.fn-ref`` markers) and are spliced back in
    as a real footnotes part after the save.
    """
    from docx import Document
    from lxml import html as LH

    doc = Document()
    frag = LH.fragment_fromstring(body_html or '<p></p>', create_parent='yr-root')

    fn_defs = _extract_footnote_defs(frag)            # strips .doc-footnotes
    fn_order = _order_footnote_refs(frag, fn_defs)    # stamps marker ids

    if frag.text and frag.text.strip():
        _add_run(doc.add_paragraph(), frag.text, {})

    for ch in frag:
        if not isinstance(ch.tag, str):
            continue
        tag = ch.tag.lower()
        if tag in _BLOCK_TAGS or tag in ('table', 'hr'):
            _emit_block(doc, ch)
        else:
            # Stray inline content at the top level → wrap in a paragraph.
            p = doc.add_paragraph()
            _place_child(p, ch, {})
        if ch.tail and ch.tail.strip():
            _add_run(doc.add_paragraph(), ch.tail, {})

    _apply_page_setup(doc, page)
    _apply_headers_footers(doc, page)
    doc.save(dest)
    if fn_order:
        try:
            _inject_footnotes(dest, fn_order)
        except Exception as e:                       # pragma: no cover - defensive
            logger.debug("footnote injection failed: %s", e)
