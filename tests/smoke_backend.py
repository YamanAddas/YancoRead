"""Quick end-to-end smoke test of the YancoRead backend.

Generates a sample file of every kind, runs each through the Flask test client,
and prints a pass/fail line per endpoint. Run with the project venv:

    venv\\Scripts\\python.exe tests\\smoke_backend.py
"""

import sys
import tempfile
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import fitz  # noqa: E402
from openpyxl import Workbook  # noqa: E402
from pptx import Presentation  # noqa: E402

from app import app, userdata  # noqa: E402

OK, FAIL = '  [ OK ]', '  [FAIL]'
results = []


def check(label, cond, extra=''):
    line = (OK if cond else FAIL) + f' {label}' + (f' — {extra}' if extra else '')
    results.append(cond)
    print(line)


def make_samples(d: Path) -> dict:
    files = {}

    # PDF (3 pages, searchable text)
    pdf = fitz.open()
    for n in range(3):
        page = pdf.new_page()
        page.insert_text((72, 90), f'YancoRead page {n + 1}\nThe quick brown fox jumps.')
    pdf.set_toc([[1, 'Start', 1], [1, 'Middle', 2]])
    p = d / 'sample.pdf'; pdf.save(p); pdf.close(); files['pdf'] = p

    # PNG image
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 200, 150))
    pix.clear_with(70)
    p = d / 'sample.png'; pix.save(p); files['image'] = p

    # CBZ comic (3 image pages)
    p = d / 'sample.cbz'
    with zipfile.ZipFile(p, 'w') as z:
        for i in range(3):
            px = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 120, 160))
            px.clear_with(40 + i * 50)
            z.writestr(f'page{i:03}.png', px.tobytes('png'))
    files['comic'] = p

    # TXT
    p = d / 'sample.txt'; p.write_text('Hello YancoRead.\nLine two.\n', encoding='utf-8'); files['text_plain'] = p

    # Markdown
    p = d / 'sample.md'
    p.write_text('# Title\n\nSome **bold** text.\n\n## Section\n\n```python\nprint("hi")\n```\n\n| a | b |\n|---|---|\n| 1 | 2 |\n', encoding='utf-8')
    files['text_md'] = p

    # Python source (code highlight path)
    p = d / 'sample.py'; p.write_text('def hello():\n    return "world"\n', encoding='utf-8'); files['text_code'] = p

    # PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Deck Title'
    slide.placeholders[1].text = 'A bullet point'
    p = d / 'sample.pptx'; prs.save(p); files['office_pptx'] = p

    # XLSX
    wb = Workbook(); ws = wb.active; ws.title = 'Data'
    ws.append(['Name', 'Score']); ws.append(['Alice', 91]); ws.append(['Bob', 88])
    p = d / 'sample.xlsx'; wb.save(p); files['office_xlsx'] = p

    return files


def main():
    tmp = Path(tempfile.mkdtemp(prefix='yancoread_smoke_'))
    files = make_samples(tmp)
    c = app.test_client()

    # health
    check('GET /health', c.get('/health').status_code == 200)

    def open_doc(path):
        r = c.post('/api/open', json={'path': str(path)})
        return r

    # PDF
    r = open_doc(files['pdf']); j = r.get_json()
    check('open PDF', r.status_code == 200 and j['doc']['kind'] == 'pdf', j.get('error', ''))
    qs = f"path={files['pdf']}"
    check('PDF page render', c.get(f'/api/page?{qs}&index=0&zoom=1').status_code == 200)
    check('PDF outline', len(c.get(f'/api/outline?{qs}').get_json()['outline']) >= 2)
    check('PDF search "fox"', len(c.get(f'/api/search?{qs}&q=fox').get_json()['results']) == 3)

    # Comic
    r = open_doc(files['comic']); j = r.get_json()
    check('open CBZ', r.status_code == 200 and j['doc']['kind'] == 'comic' and j['doc']['meta']['page_count'] == 3, j.get('error', ''))
    check('comic page', c.get(f"/api/comic-page?path={files['comic']}&index=1").status_code == 200)

    # Image
    r = open_doc(files['image']); j = r.get_json()
    check('open PNG', r.status_code == 200 and j['doc']['kind'] == 'image', j.get('error', ''))
    check('image serve', c.get(f"/api/image?path={files['image']}").status_code == 200)

    # Text variants
    for key, label, mode in [('text_plain', 'TXT', 'plain'), ('text_md', 'Markdown', 'markdown'), ('text_code', 'code', 'code')]:
        r = open_doc(files[key]); j = r.get_json()
        check(f'open {label}', r.status_code == 200 and j['doc']['kind'] == 'text', j.get('error', ''))
        tr = c.get(f"/api/text?path={files[key]}").get_json()
        check(f'{label} render (mode={mode})', tr.get('mode') == mode and bool(tr.get('html')), str(tr.get('mode')))

    # Office
    for key, label in [('office_pptx', 'PPTX'), ('office_xlsx', 'XLSX')]:
        r = open_doc(files[key]); j = r.get_json()
        check(f'open {label}', r.status_code == 200 and j['doc']['kind'] == 'office', j.get('error', ''))
        orr = c.get(f"/api/office?path={files[key]}").get_json()
        check(f'{label} render', bool(orr.get('html')) and len(orr.get('outline', [])) >= 1)

    # cleanup recents created during the smoke run
    userdata.clear_recent()

    print('\n' + ('ALL PASSED' if all(results) else f'{results.count(False)} FAILED') +
          f'  ({results.count(True)}/{len(results)})')
    sys.exit(0 if all(results) else 1)


if __name__ == '__main__':
    main()
