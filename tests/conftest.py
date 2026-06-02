"""Shared pytest fixtures for YancoRead: generated sample files + Flask client."""
import io
import sys
import zipfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))


def _png(w, h, val):
    import cv2
    import numpy as np
    img = np.full((h, w, 3), val, np.uint8)
    return cv2.imencode('.png', img)[1].tobytes()


def _panel_png():
    """A 2x2-grid comic page (white gutters) → 4 detectable panels."""
    import cv2
    import numpy as np
    W, H, g = 900, 1300, 26
    img = np.full((H, W, 3), 250, np.uint8)
    rects = [(0, 0, .5, .5), (.5, 0, 1, .5), (0, .5, .5, 1), (.5, .5, 1, 1)]
    cols = [(60, 120, 200), (200, 120, 60), (90, 180, 110), (170, 90, 180)]
    for (x0, y0, x1, y1), c in zip(rects, cols):
        a = (int(x0 * W) + g, int(y0 * H) + g)
        b = (int(x1 * W) - g, int(y1 * H) - g)
        cv2.rectangle(img, a, b, c, -1)
        cv2.rectangle(img, a, b, (20, 20, 20), 5)
    return cv2.imencode('.png', img)[1].tobytes()


@pytest.fixture(scope='session')
def samples(tmp_path_factory):
    import fitz
    from openpyxl import Workbook
    from pptx import Presentation

    d = tmp_path_factory.mktemp('yr_samples')
    files = {}

    pdf = fitz.open()
    for n in range(3):
        pg = pdf.new_page()
        pg.insert_text((72, 90), f'YancoRead page {n + 1}\nThe quick brown fox.')
    pdf.set_toc([[1, 'Start', 1], [1, 'Middle', 2]])
    files['pdf'] = d / 'a.pdf'; pdf.save(files['pdf']); pdf.close()

    files['png'] = d / 'a.png'; files['png'].write_bytes(_png(200, 150, 80))

    files['cbz'] = d / 'a.cbz'
    with zipfile.ZipFile(files['cbz'], 'w') as z:
        for i in range(3):
            z.writestr(f'p{i:02}.png', _png(120, 160, 40 + i * 40))

    files['cbz_panels'] = d / 'panels.cbz'
    with zipfile.ZipFile(files['cbz_panels'], 'w') as z:
        z.writestr('p00.png', _panel_png())

    files['cbz_manga'] = d / 'manga.cbz'
    with zipfile.ZipFile(files['cbz_manga'], 'w') as z:
        z.writestr('ComicInfo.xml', '<?xml version="1.0"?><ComicInfo><Manga>YesAndRightToLeft</Manga></ComicInfo>')
        z.writestr('p00.png', _png(120, 160, 50))

    files['txt'] = d / 'a.txt'; files['txt'].write_text('Hello YancoRead.\nLine two.\n', encoding='utf-8')
    files['md'] = d / 'a.md'; files['md'].write_text('# Title\n\n**bold**\n\n## Sec\n\n| a | b |\n|---|---|\n| 1 | 2 |\n', encoding='utf-8')
    files['py'] = d / 'a.py'; files['py'].write_text('def hi():\n    return 1\n', encoding='utf-8')

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = 'Deck'; s.placeholders[1].text = 'Bullet'
    files['pptx'] = d / 'a.pptx'; prs.save(files['pptx'])

    wb = Workbook(); ws = wb.active; ws.title = 'Data'
    ws.append(['Name', 'Score']); ws.append(['Alice', 91])
    files['xlsx'] = d / 'a.xlsx'; wb.save(files['xlsx'])

    return files


@pytest.fixture
def client(tmp_path):
    """Flask test client with user data isolated to a temp file."""
    import app as app_module
    from userdata import UserData
    app_module.app.config['TESTING'] = True
    original = app_module.userdata
    app_module.userdata = UserData(tmp_path / 'userdata.json')
    try:
        c = app_module.app.test_client()
        # Send the per-session API token on every request so the write-guard
        # doesn't 403 normal tests (mirrors how the real frontend behaves).
        c.environ_base['HTTP_X_YR_TOKEN'] = app_module._API_TOKEN
        yield c
    finally:
        app_module.userdata = original
