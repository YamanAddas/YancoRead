"""Generate persistent sample files for manual / preview testing.

    venv\\Scripts\\python.exe tests\\make_samples.py [out_dir]
"""
import sys
import zipfile
from pathlib import Path

import fitz
from openpyxl import Workbook
from pptx import Presentation

out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent / 'samples'
out.mkdir(parents=True, exist_ok=True)

# PDF
pdf = fitz.open()
for n in range(6):
    pg = pdf.new_page()
    pg.insert_text((72, 100), f'YancoRead — sample PDF\n\nPage {n + 1} of 6\n'
                              'The quick brown fox jumps over the lazy dog.', fontsize=16)
pdf.set_toc([[1, 'Cover', 1], [1, 'Chapter One', 2], [1, 'Chapter Two', 4]])
pdf.save(out / 'sample.pdf'); pdf.close()

# PNG
px = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 600, 400)); px.clear_with(60); px.save(out / 'sample.png')

# CBZ
with zipfile.ZipFile(out / 'sample.cbz', 'w') as z:
    for i in range(5):
        p = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 800, 1200)); p.clear_with(30 + i * 35)
        z.writestr(f'page{i:03}.png', p.tobytes('png'))

# EPUB (reflowable — exercises the fitz eBook path: search, outline, relayout)
_epub_chap = ('<?xml version="1.0" encoding="utf-8"?>\n'
              '<!DOCTYPE html><html xmlns="http://www.w3.org/1999/xhtml"><head><title>Ch1</title></head>'
              '<body><h1>Chapter One</h1>'
              '<p>The quick brown fox jumps over the lazy dog.</p>'
              '<p>YancoRead eBook sample paragraph for searching.</p></body></html>')
_epub_opf = ('<?xml version="1.0"?>\n'
             '<package xmlns="http://www.idpf.org/2007/opf" unique-identifier="id" version="2.0">'
             '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>YancoRead Sample EPUB</dc:title>'
             '<dc:creator>YancoRead</dc:creator><dc:identifier id="id">yr-sample-1</dc:identifier>'
             '<dc:language>en</dc:language></metadata>'
             '<manifest><item id="c1" href="chapter1.xhtml" media-type="application/xhtml+xml"/>'
             '<item id="ncx" href="toc.ncx" media-type="application/x-dtbncx+xml"/></manifest>'
             '<spine toc="ncx"><itemref idref="c1"/></spine></package>')
_epub_ncx = ('<?xml version="1.0"?>\n'
             '<ncx xmlns="http://www.daisy.org/z3986/2005/ncx/" version="2005-1"><head/>'
             '<docTitle><text>YancoRead Sample EPUB</text></docTitle>'
             '<navMap><navPoint id="n1" playOrder="1"><navLabel><text>Chapter One</text></navLabel>'
             '<content src="chapter1.xhtml"/></navPoint></navMap></ncx>')
_epub_container = ('<?xml version="1.0"?>\n'
                   '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
                   '<rootfiles><rootfile full-path="OEBPS/content.opf" '
                   'media-type="application/oebps-package+xml"/></rootfiles></container>')
with zipfile.ZipFile(out / 'sample.epub', 'w') as z:
    z.writestr('mimetype', 'application/epub+zip', compress_type=zipfile.ZIP_STORED)
    z.writestr('META-INF/container.xml', _epub_container, compress_type=zipfile.ZIP_DEFLATED)
    z.writestr('OEBPS/content.opf', _epub_opf, compress_type=zipfile.ZIP_DEFLATED)
    z.writestr('OEBPS/toc.ncx', _epub_ncx, compress_type=zipfile.ZIP_DEFLATED)
    z.writestr('OEBPS/chapter1.xhtml', _epub_chap, compress_type=zipfile.ZIP_DEFLATED)

# TXT / MD / code
(out / 'sample.txt').write_text('YancoRead plain text sample.\n' + 'Lorem ipsum dolor sit amet.\n' * 40, encoding='utf-8')
(out / 'sample.md').write_text(
    '# YancoRead Markdown\n\nA **universal** reader.\n\n## Features\n\n'
    '- PDF\n- Comics\n- eBooks\n\n## Code\n\n```python\ndef hi():\n    print("hi")\n```\n\n'
    '## Table\n\n| Format | Tool |\n|---|---|\n| PDF | zoom |\n| Comic | spread |\n', encoding='utf-8')
(out / 'sample.py').write_text(
    'import sys\n\n\ndef main():\n    """Say hello."""\n    print("Hello from YancoRead")\n\n\n'
    'if __name__ == "__main__":\n    main()\n', encoding='utf-8')

# PPTX
prs = Presentation()
for title, body in [('YancoRead', 'Universal document reader'), ('Slide Two', 'Adaptive tools per format')]:
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    s.placeholders[1].text = body
prs.save(out / 'sample.pptx')

# XLSX
wb = Workbook(); ws = wb.active; ws.title = 'Scores'
ws.append(['Name', 'Score', 'Grade'])
for row in [['Alice', 91, 'A'], ['Bob', 88, 'B'], ['Cara', 95, 'A']]:
    ws.append(row)
wb.save(out / 'sample.xlsx')

print('Samples written to', out)
for f in sorted(out.iterdir()):
    print('  ', f.name)
