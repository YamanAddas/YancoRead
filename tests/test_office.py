"""Office rendering-mode selection (native-only: docx/pptx/xlsx via HTML)."""

import re
import zipfile

import docx

import app
from renderers import officedoc


def _make_review_docx(path):
    """A .docx with a comment (on para 1) and a tracked insertion + deletion
    (in para 2), built by injecting raw OXML for the tracked changes."""
    from docx import Document
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    d = Document()
    p1 = d.add_paragraph('The quick brown fox')
    d.add_comment(p1.runs, text='Tighten this sentence', author='Alice', initials='A')

    p2 = d.add_paragraph()
    ins = OxmlElement('w:ins')
    ins.set(qn('w:author'), 'Bob'); ins.set(qn('w:date'), '2024-05-01T10:00:00Z'); ins.set(qn('w:id'), '11')
    r = OxmlElement('w:r'); t = OxmlElement('w:t'); t.text = 'added clause'; r.append(t); ins.append(r)
    p2._p.append(ins)
    dele = OxmlElement('w:del')
    dele.set(qn('w:author'), 'Bob'); dele.set(qn('w:date'), '2024-05-01T10:01:00Z'); dele.set(qn('w:id'), '12')
    r2 = OxmlElement('w:r'); dt = OxmlElement('w:delText'); dt.text = 'old phrase'; r2.append(dt); dele.append(r2)
    p2._p.append(dele)
    d.save(str(path))


def test_docx_review_extraction(tmp_path):
    """Tracked changes + comments are extracted with authors and anchored text."""
    f = tmp_path / 'review.docx'; _make_review_docx(f)
    assert officedoc._docx_has_review(str(f)) is True

    out = officedoc.to_html(str(f))
    rv = out['review']
    types = {(c['type'], c['author'], c['text']) for c in rv['changes']}
    assert ('ins', 'Bob', 'added clause') in types
    assert ('del', 'Bob', 'old phrase') in types
    cm = rv['comments'][0]
    assert cm['author'] == 'Alice' and cm['text'] == 'Tighten this sentence'
    assert 'quick brown fox' in cm['quote']      # anchored text captured
    assert set(rv['authors']) == {'Alice', 'Bob'}


def test_docx_markup_body(tmp_path):
    """The markup body shows insertions and deletions inline."""
    f = tmp_path / 'review.docx'; _make_review_docx(f)
    html = officedoc.to_html(str(f))['markupHtml']
    assert 'class="trk-ins"' in html and 'added clause</ins>' in html
    assert 'class="trk-del"' in html and 'old phrase</del>' in html
    assert 'data-author="Bob"' in html
    assert 'cmt-anchor' in html               # comment anchor marker present


def test_docx_compare_no_backup(tmp_path):
    """No .bak yet → compare reports backup:false (and never raises)."""
    from docx import Document
    f = tmp_path / 'doc.docx'
    d = Document(); d.add_paragraph('Original line.'); d.save(str(f))
    assert officedoc._docx_compare(str(f)) == {'backup': False}
    assert officedoc.to_html(str(f))['hasBackup'] is False


def test_docx_compare_redline(tmp_path):
    """A redline marks deleted/inserted/edited paragraphs against the backup."""
    from docx import Document
    f = tmp_path / 'doc.docx'
    bak = tmp_path / 'doc.docx.bak'

    # backup = the "before" version
    d0 = Document()
    d0.add_paragraph('Keep this line.')
    d0.add_paragraph('Delete this whole line.')
    d0.add_paragraph('The price is ten dollars.')
    d0.save(str(bak))
    # current = the "after" version
    d1 = Document()
    d1.add_paragraph('Keep this line.')
    d1.add_paragraph('The price is twenty dollars.')   # word-level edit
    d1.add_paragraph('A brand new line.')              # inserted
    d1.save(str(f))

    out = officedoc._docx_compare(str(f))
    assert out['backup'] is True and out['changed'] >= 2
    h = out['html']
    assert 'Delete this whole line.</del>' in h          # deleted paragraph
    assert 'A brand new line.</ins>' in h                # inserted paragraph
    assert '<del class="trk-del">ten</del>' in h or 'ten</del>' in h   # word-level delete
    assert 'twenty' in h and 'trk-ins' in h              # word-level insert
    assert '<p>Keep this line.</p>' in h                 # unchanged paragraph stays plain
    assert officedoc.to_html(str(f))['hasBackup'] is True


def test_html_to_markdown():
    """Core HTML constructs convert to clean Markdown."""
    html = (
        '<article class="doc-page docx">'
        '<h1>Title</h1>'
        '<h2>Section</h2>'
        '<p>Plain with <strong>bold</strong>, <em>italic</em> and a '
        '<a href="https://x.test">link</a>.</p>'
        '<ul><li>one</li><li>two<ul><li>nested</li></ul></li></ul>'
        '<ol><li>first</li><li>second</li></ol>'
        '<blockquote><p>quoted</p></blockquote>'
        '<table><tr><th>A</th><th>B</th></tr><tr><td>1</td><td>2</td></tr></table>'
        '</article>')
    md = officedoc.html_to_markdown(html)
    assert '# Title' in md and '## Section' in md
    assert '**bold**' in md and '*italic*' in md
    assert '[link](https://x.test)' in md
    assert '- one' in md and '  - nested' in md       # nested list indent
    assert '1. first' in md and '2. second' in md
    assert '> quoted' in md
    assert '| A | B |' in md and '| --- | --- |' in md and '| 1 | 2 |' in md


def test_html_to_markdown_skips_embedded_images():
    md = officedoc.html_to_markdown('<article><p><img alt="logo" src="data:image/png;base64,AAAA"></p></article>')
    assert '![logo](embedded image)' in md and 'base64' not in md


def test_office_export_writes_markdown(client, tmp_path):
    """/api/office/export re-renders the .docx from disk and writes Markdown."""
    from docx import Document
    f = tmp_path / 'doc.docx'
    d = Document(); d.add_heading('Report', 0); d.add_paragraph('Hello world.'); d.save(str(f))
    out = tmp_path / 'out.md'
    r = client.post('/api/office/export', json={'path': str(f), 'format': 'md', 'target': str(out)})
    j = r.get_json()
    assert j.get('ok') is True
    text = out.read_text(encoding='utf-8')
    assert 'Hello world.' in text and 'Report' in text


def test_office_export_rejects_non_docx(client, tmp_path):
    bad = tmp_path / 'x.txt'; bad.write_text('hi', encoding='utf-8')
    r = client.post('/api/office/export', json={'path': str(bad), 'format': 'md', 'target': str(tmp_path / 'o.md')})
    assert r.status_code >= 400


def test_docx_accept_changes(tmp_path):
    """Accept keeps insertions, drops deletions → a clean .docx with no
    remaining tracked-change marks."""
    from docx import Document
    src = tmp_path / 'rev.docx'; _make_review_docx(src)   # ins 'added clause', del 'old phrase'
    out = tmp_path / 'accepted.docx'
    res = officedoc.accept_reject_changes(str(src), str(out), 'accept')
    assert res['ok'] and res['changed'] >= 2

    text = '\n'.join(p.text for p in Document(str(out)).paragraphs)
    assert 'added clause' in text          # insertion kept
    assert 'old phrase' not in text        # deletion dropped
    with zipfile.ZipFile(str(out)) as z:
        dx = z.read('word/document.xml').decode('utf-8')
    assert re.search(r'<w:(ins|del)[ >]', dx) is None       # no change marks remain


def test_docx_reject_changes(tmp_path):
    """Reject drops insertions, restores deletions."""
    from docx import Document
    src = tmp_path / 'rev.docx'; _make_review_docx(src)
    out = tmp_path / 'rejected.docx'
    officedoc.accept_reject_changes(str(src), str(out), 'reject')

    text = '\n'.join(p.text for p in Document(str(out)).paragraphs)
    assert 'added clause' not in text      # insertion dropped
    assert 'old phrase' in text            # deletion restored as live text
    with zipfile.ZipFile(str(out)) as z:
        dx = z.read('word/document.xml').decode('utf-8')
    assert re.search(r'<w:(ins|del)[ >]', dx) is None
    assert 'delText' not in dx             # restored text is live <w:t>, not <w:delText>


def test_docx_accept_changes_preserves_original(tmp_path):
    """The source file is never modified by accept/reject."""
    src = tmp_path / 'rev.docx'; _make_review_docx(src)
    before = src.read_bytes()
    officedoc.accept_reject_changes(str(src), str(tmp_path / 'o.docx'), 'accept')
    assert src.read_bytes() == before


def test_accept_changes_endpoint(client, tmp_path):
    """/api/office/accept-changes writes a new .docx and validates mode."""
    src = tmp_path / 'rev.docx'; _make_review_docx(src)
    out = tmp_path / 'done.docx'
    r = client.post('/api/office/accept-changes',
                    json={'path': str(src), 'mode': 'accept', 'target': str(out)})
    j = r.get_json()
    assert j.get('ok') is True and out.is_file()
    bad = client.post('/api/office/accept-changes',
                      json={'path': str(src), 'mode': 'sideways', 'target': str(out)})
    assert bad.status_code >= 400


def test_docx_without_review_has_no_review_key(tmp_path):
    """A plain document carries no review payload (zero cost)."""
    from docx import Document
    d = Document(); d.add_paragraph('Just plain text.'); d.save(str(tmp_path / 'p.docx'))
    out = officedoc.to_html(str(tmp_path / 'p.docx'))
    assert 'review' not in out and 'markupHtml' not in out
    assert officedoc._docx_has_review(str(tmp_path / 'p.docx')) is False


def test_office_meta_native_is_flow():
    """DOCX opens via the lightweight native HTML (flow) path."""
    meta = app._office_meta('whatever.docx', '.docx')
    assert meta['render'] == 'flow'


def test_office_meta_pptx_is_slides():
    """PPTX gets the dedicated one-slide-at-a-time deck viewer."""
    for ext in ('.pptx', '.PPTX'):
        meta = app._office_meta('whatever' + ext, ext)
        assert meta['render'] == 'slides', ext


def test_office_meta_xlsx_is_sheet():
    """XLSX gets the structured sticky-grid spreadsheet viewer."""
    for ext in ('.xlsx', '.XLSX'):
        meta = app._office_meta('whatever' + ext, ext)
        assert meta['render'] == 'sheet', ext


def test_xlsx_structured_output(tmp_path):
    """_xlsx_to_html returns per-sheet structure: sparse cells, merges,
    freeze panes, and column/row sizes."""
    import datetime
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active; ws.title = 'Sales'
    ws['A1'] = 'Region'; ws['B1'] = 'Q1'
    ws['A2'] = 'EU'; ws['B2'] = 1200.5
    ws['A3'] = 'Date'; ws['B3'] = datetime.datetime(2024, 1, 15)
    ws.merge_cells('A5:C5'); ws['A5'] = 'Footer'
    ws.freeze_panes = 'B2'
    ws.column_dimensions['A'].width = 18
    ws.row_dimensions[1].height = 24
    wb.create_sheet('Empty')
    f = tmp_path / 's.xlsx'; wb.save(str(f))

    out = officedoc.to_html(str(f))
    assert [s['name'] for s in out['sheets']] == ['Sales', 'Empty']
    sales = out['sheets'][0]
    # sparse cells: only non-empty are present (covered merge cells excluded)
    by_rc = {(c['r'], c['c']): c['v'] for c in sales['cells']}
    assert by_rc[(1, 1)] == 'Region' and by_rc[(2, 2)] == 1200.5
    # datetime(2024,1,15) → Excel serial number (for SSF formatting on the front)
    assert by_rc[(3, 2)] == 45306.0
    assert (5, 2) not in by_rc                        # covered by the merge
    assert {'r': 5, 'c': 1, 'rs': 1, 'cs': 3} in sales['merged']
    assert sales['freeze'] == 'B2'
    assert sales['colWidths'].get(1) and sales['rowHeights'].get(1)


def test_xlsx_cell_styles(tmp_path):
    """Cells carry compact style flags; theme/indexed colors are skipped."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.styles.colors import Color

    wb = Workbook(); ws = wb.active
    ws['A1'] = 'Hdr'
    ws['A1'].font = Font(bold=True, italic=True, color='FFCC0000')
    ws['A1'].fill = PatternFill('solid', fgColor='FFEEF3F8')
    ws['A1'].alignment = Alignment(horizontal='center')
    ws['A1'].border = Border(bottom=Side(style='thin'), right=Side(style='medium'))
    ws['A2'] = 'themed'
    ws['A2'].font = Font(color=Color(theme=1, type='theme'))   # must be skipped
    f = tmp_path / 'styled.xlsx'; wb.save(str(f))

    cells = {(c['r'], c['c']): c for c in officedoc.to_html(str(f))['sheets'][0]['cells']}
    s = cells[(1, 1)]['s']
    assert s['b'] == 1 and s['i'] == 1
    assert s['fc'] == '#cc0000' and s['bg'] == '#eef3f8'
    assert s['a'] == 'c'
    assert 'b' in s['bd'] and 'r' in s['bd'] and 't' not in s['bd']
    # themed font color is dropped → no 'fc' (cell may have no style at all)
    assert 'fc' not in cells[(2, 1)].get('s', {})


def test_xlsx_formula_exposure(tmp_path):
    """Formula cells carry their formula string in `f`; a formula with no cached
    value still appears (driven by the formula pass), not dropped as empty."""
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active
    ws['A1'] = 10; ws['A2'] = 20
    ws['A3'] = '=SUM(A1:A2)'      # openpyxl writes no cached value
    ws['B1'] = 'plain'
    f = tmp_path / 'fm.xlsx'; wb.save(str(f))

    cells = {(c['r'], c['c']): c for c in officedoc.to_html(str(f))['sheets'][0]['cells']}
    assert cells[(3, 1)]['f'] == '=SUM(A1:A2)'   # formula present
    assert cells[(3, 1)]['v'] is None            # no cached value
    assert 'f' not in cells[(1, 1)]              # literal has no formula
    assert cells[(2, 1)]['v'] == 20


def test_xlsx_caps_huge_sheet(tmp_path):
    """Rows/cols past the cap are dropped and flagged truncated."""
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active
    ws.cell(row=officedoc._XLSX_MAX_ROWS + 50, column=1, value='past-cap')
    f = tmp_path / 'big.xlsx'; wb.save(str(f))
    sheet = officedoc.to_html(str(f))['sheets'][0]
    assert sheet['truncated'] is True
    assert sheet['rows'] == officedoc._XLSX_MAX_ROWS


def test_pptx_render_emits_slide_geometry(tmp_path):
    """_pptx_to_html returns per-slide sections + stage size + count."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9 widescreen (13.333" × 7.5")
    prs.slide_height = Emu(6858000)
    for title, body in [('One', 'first'), ('Two', 'second'), ('Three', 'third')]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        s.placeholders[1].text = body
    f = tmp_path / 'deck.pptx'
    prs.save(str(f))

    out = officedoc.to_html(str(f))
    assert out['slide_count'] == 3
    assert len(out['outline']) == 3
    assert out['html'].count('class="slide"') == 3
    # 12192000 / 9525 = 1280, 6858000 / 9525 = 720  → exact 16:9.
    assert out['slide_size'] == {'width': 1280, 'height': 720}
    # Slide anchors are stable ids the frontend splits on.
    assert 'id="slide-1"' in out['html'] and 'id="slide-3"' in out['html']


def test_pptx_positions_shapes_at_native_px(tmp_path):
    """Shapes are emitted as absolutely-positioned divs at native EMU/9525 px,
    with run styling (size→px, bold, color)."""
    from pptx import Presentation
    from pptx.util import Emu, Pt, Inches
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    s = prs.slides.add_slide(prs.slide_layouts[6])     # blank
    tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = 'Hi'; r.font.size = Pt(36); r.font.bold = True
    r.font.color.rgb = RGBColor(0x20, 0x40, 0x80)
    f = tmp_path / 'pos.pptx'; prs.save(str(f))

    html = officedoc.to_html(str(f))['html']
    # 1in=96px, 0.5in=48px, 5in=480px, 1.5in=144px
    assert 'left:96.0px;top:48.0px;width:480.0px;height:144.0px' in html
    assert 'class="sl-shape"' in html and 'class="sl-text"' in html
    # 36pt → 48px; bold; explicit RGB.
    assert 'font-size:48.0px' in html and 'font-weight:700' in html
    assert 'color:#204080' in html


def test_pptx_speaker_notes(tmp_path):
    """notes[] carries per-slide speaker notes; slides without notes give ''
    (and accessing them must not fabricate a notes slide)."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    s1.notes_slide.notes_text_frame.text = 'Mention the Q3 numbers.'
    prs.slides.add_slide(prs.slide_layouts[6])      # no notes
    f = tmp_path / 'notes.pptx'; prs.save(str(f))

    out = officedoc.to_html(str(f))
    assert out['notes'] == ['Mention the Q3 numbers.', '']
    assert len(out['notes']) == out['slide_count']


def test_pptx_table_and_group_and_bg(tmp_path):
    """Tables render as <table class=sl-table>, grouped shapes are mapped onto
    the slide via the group transform, and a solid slide background is emitted."""
    from pptx import Presentation
    from pptx.util import Emu, Inches
    from pptx.enum.dml import MSO_FILL_TYPE
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Emu(12192000); prs.slide_height = Emu(6858000)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # solid slide background
    bg = s.background.fill
    bg.solid(); bg.fore_color.rgb = RGBColor(0x10, 0x20, 0x30)

    t = s.shapes.add_table(2, 2, Inches(6), Inches(1), Inches(4), Inches(2)).table
    t.cell(0, 0).text = 'A'; t.cell(1, 1).text = 'D'

    grp = s.shapes.add_group_shape()
    inner = grp.shapes.add_textbox(Inches(2), Inches(3), Inches(3), Inches(1))
    inner.text_frame.text = 'grouped'

    f = tmp_path / 'tg.pptx'; prs.save(str(f))
    html = officedoc.to_html(str(f))['html']

    assert 'sl-tablebox' in html and 'class="sl-table"' in html
    assert '>A</td>' in html and '>D</td>' in html
    assert 'sl-bg' in html and 'background:#102030' in html
    assert 'sl-bg-dark' in html   # #102030 is dark → light default text
    # grouped textbox at child (2in,3in) → slide px (192,288) via identity group xfrm
    assert 'left:192.0px;top:288.0px' in html
    assert 'grouped' in html


def test_office_meta_legacy_is_unsupported():
    """Legacy / OpenDocument formats can't be opened natively and report so."""
    for ext in ('.doc', '.ppt', '.xls', '.rtf', '.odt', '.odp', '.ods'):
        meta = app._office_meta('whatever' + ext, ext)
        assert meta['render'] == 'unsupported', ext
        assert meta['reason'] == 'unsupported_format'
        assert meta['ext'] == ext


def test_office_meta_handles_uppercase_ext():
    meta = app._office_meta('WHATEVER.DOCX', '.DOCX')
    assert meta['render'] == 'flow'


# ── editor write-back (HTML → docx) ────────────────────────────────────────────
RICH_HTML = (
    '<h1>Title</h1>'
    '<p>Plain <b>bold</b> <i>italic</i> <u>underline</u> <s>strike</s>.</p>'
    '<p style="text-align:center">'
    '<span style="color:rgb(224,49,49);font-size:20px">red big centred</span></p>'
    '<p style="line-height:2">'
    '<span style="font-family:Georgia;background-color:#ffe066">georgia hl</span></p>'
    '<ul><li>a</li><li>b</li></ul>'
    '<ol><li>one</li><li>two</li></ol>'
    '<blockquote>quoted</blockquote>'
    '<table><tr><th>H</th></tr><tr><td>c</td></tr></table>'
)


def _by_text(d, text):
    for p in d.paragraphs:
        if p.text.strip() == text:
            return p
    raise AssertionError(f'paragraph {text!r} not found')


def test_html_to_docx_preserves_formatting(tmp_path):
    dest = tmp_path / 'out.docx'
    officedoc.html_to_docx(RICH_HTML, str(dest))
    d = docx.Document(str(dest))

    assert _by_text(d, 'Title').style.name == 'Heading 1'

    runs = {r.text: r for p in d.paragraphs for r in p.runs}
    assert runs['bold'].bold is True
    assert runs['italic'].italic is True
    assert runs['underline'].underline is True
    assert runs['strike'].font.strike is True

    from docx.enum.text import WD_ALIGN_PARAGRAPH
    centred = _by_text(d, 'red big centred')
    assert centred.alignment == WD_ALIGN_PARAGRAPH.CENTER
    big = centred.runs[0]
    assert str(big.font.color.rgb) == 'E03131'
    assert abs(big.font.size.pt - 15.0) < 0.5      # 20px → 15pt

    geo = _by_text(d, 'georgia hl')
    assert geo.paragraph_format.line_spacing == 2.0
    assert geo.runs[0].font.name == 'Georgia'

    styles = [p.style.name for p in d.paragraphs]
    assert 'List Bullet' in styles
    assert 'List Number' in styles
    assert 'Quote' in styles

    assert len(d.tables) == 1
    assert d.tables[0].cell(1, 0).text == 'c'


def test_html_to_docx_preserves_hyperlink(tmp_path):
    """An <a href> becomes a real Word hyperlink with an external relationship."""
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    dest = tmp_path / 'link.docx'
    officedoc.html_to_docx(
        '<p>see <a href="https://example.com/page">this link</a> now</p>', str(dest))
    d = docx.Document(str(dest))

    targets = [r.target_ref for r in d.part.rels.values() if r.reltype == RT.HYPERLINK]
    assert 'https://example.com/page' in targets

    xml = d.paragraphs[0]._p.xml
    assert 'w:hyperlink' in xml          # a genuine hyperlink element…
    assert 'this link' in xml            # …wrapping the visible text


def test_html_to_docx_rejects_javascript_href(tmp_path):
    """Unsafe link schemes don't get a relationship (but the text survives)."""
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    dest = tmp_path / 'js.docx'
    officedoc.html_to_docx('<p><a href="javascript:alert(1)">click</a></p>', str(dest))
    d = docx.Document(str(dest))

    targets = [r.target_ref for r in d.part.rels.values() if r.reltype == RT.HYPERLINK]
    assert not targets
    assert any('click' in p.text for p in d.paragraphs)


def test_office_save_overwrite_makes_backup(tmp_path):
    src = tmp_path / 'doc.docx'
    docx.Document().save(str(src))           # a real, openable docx

    client = app.app.test_client()
    resp = client.post('/api/office/save', json={
        'path': str(src), 'mode': 'overwrite',
        'html': '<h1>Hello</h1><p>world</p>',
    })
    assert resp.status_code == 200
    data = resp.get_json()
    assert data['ok'] is True
    assert data['backup'] and (tmp_path / 'doc.docx.bak').is_file()

    saved = docx.Document(str(src))
    assert _by_text(saved, 'Hello').style.name == 'Heading 1'
    assert any(p.text == 'world' for p in saved.paragraphs)


def test_office_save_as_writes_new_file_and_coerces_ext(tmp_path):
    target = tmp_path / 'copy.txt'           # wrong suffix on purpose
    client = app.app.test_client()
    resp = client.post('/api/office/save', json={
        'mode': 'saveas', 'target': str(target),
        'html': '<p>fresh copy</p>',
    })
    assert resp.status_code == 200
    data = resp.get_json()
    assert data['backup'] is None
    out = tmp_path / 'copy.docx'             # coerced to .docx
    assert out.is_file()
    assert any(p.text == 'fresh copy' for p in docx.Document(str(out)).paragraphs)


def test_office_save_overwrite_rejects_non_docx(tmp_path):
    bad = tmp_path / 'note.txt'
    bad.write_text('hi', encoding='utf-8')
    client = app.app.test_client()
    resp = client.post('/api/office/save', json={
        'path': str(bad), 'mode': 'overwrite', 'html': '<p>x</p>',
    })
    assert resp.status_code == 400
    assert 'docx' in resp.get_json()['error'].lower()


# ── round-trip fidelity detection (gate lossy overwrites) ──────────────────────
def test_fidelity_plain_doc_is_clean(tmp_path):
    """An ordinary text/heading/table doc round-trips, so nothing is flagged."""
    src = tmp_path / 'plain.docx'
    d = docx.Document()
    d.add_heading('Title', level=1)
    d.add_paragraph('Just some body text.')
    d.add_table(rows=2, cols=2)
    d.save(str(src))

    fid = officedoc.detect_docx_fidelity(str(src))
    assert fid['lossy'] is False
    assert fid['features'] == []


def test_fidelity_plain_header_round_trips_clean(tmp_path):
    """A single plain-text primary header now round-trips (Tier 4b), so it is
    no longer flagged as lossy."""
    src = tmp_path / 'header.docx'
    d = docx.Document()
    d.add_paragraph('body')
    d.sections[0].header.paragraphs[0].text = 'Confidential — Draft'
    d.save(str(src))

    fid = officedoc.detect_docx_fidelity(str(src))
    assert 'Headers & footers' not in fid['features']


def test_fidelity_flags_first_page_header(tmp_path):
    """A distinct first-page header is a variant the rebuild can't keep, so the
    file is still steered to Save As."""
    src = tmp_path / 'firsthdr.docx'
    d = docx.Document()
    d.add_paragraph('body')
    s = d.sections[0]
    s.different_first_page_header_footer = True
    s.first_page_header.paragraphs[0].text = 'first page only'
    s.header.paragraphs[0].text = 'every page'
    d.save(str(src))

    fid = officedoc.detect_docx_fidelity(str(src))
    assert fid['lossy'] is True
    assert 'Headers & footers' in fid['features']


def test_fidelity_flags_header_with_table(tmp_path):
    """A header holding a table (multi-column layout) is rich → flagged."""
    from docx.shared import Inches
    src = tmp_path / 'tblhdr.docx'
    d = docx.Document()
    d.add_paragraph('body')
    hdr = d.sections[0].header
    hdr.paragraphs[0].text = 'left'
    hdr.add_table(rows=1, cols=2, width=Inches(6))
    d.save(str(src))

    fid = officedoc.detect_docx_fidelity(str(src))
    assert 'Headers & footers' in fid['features']


def test_fidelity_flags_multiple_sections(tmp_path):
    """A second section (the body then holds >1 sectPr) is unpreservable."""
    from docx.enum.section import WD_SECTION
    src = tmp_path / 'sections.docx'
    d = docx.Document()
    d.add_paragraph('first section')
    d.add_section(WD_SECTION.NEW_PAGE)
    d.add_paragraph('second section')
    d.save(str(src))

    fid = officedoc.detect_docx_fidelity(str(src))
    assert fid['lossy'] is True
    assert 'Multiple sections' in fid['features']


def test_fidelity_never_raises_on_bad_input(tmp_path):
    """A non-zip / missing file reports 'nothing lost' instead of blowing up."""
    bad = tmp_path / 'not-a-docx.docx'
    bad.write_text('definitely not a zip', encoding='utf-8')
    assert officedoc.detect_docx_fidelity(str(bad)) == {'lossy': False, 'features': []}
    assert officedoc.detect_docx_fidelity(str(tmp_path / 'missing.docx')) == {
        'lossy': False, 'features': []}


# ── table editing round-trip (merge / header / column width) ───────────────────
def test_table_colspan_merges_horizontally(tmp_path):
    dest = tmp_path / 't.docx'
    officedoc.html_to_docx(
        '<table><tr><td colspan="2">wide</td></tr>'
        '<tr><td>a</td><td>b</td></tr></table>', str(dest))
    t = docx.Document(str(dest)).tables[0]
    # The two top grid cells resolve to the same underlying <w:tc>.
    assert t.cell(0, 0)._tc is t.cell(0, 1)._tc
    assert t.cell(0, 0).text == 'wide'
    assert t.cell(1, 0).text == 'a' and t.cell(1, 1).text == 'b'


def test_table_rowspan_merges_vertically(tmp_path):
    dest = tmp_path / 't.docx'
    officedoc.html_to_docx(
        '<table><tr><td rowspan="2">tall</td><td>x</td></tr>'
        '<tr><td>y</td></tr></table>', str(dest))
    t = docx.Document(str(dest)).tables[0]
    assert t.cell(0, 0)._tc is t.cell(1, 0)._tc
    assert t.cell(0, 0).text == 'tall'
    assert t.cell(0, 1).text == 'x' and t.cell(1, 1).text == 'y'


def test_table_th_row_is_bold_and_repeats(tmp_path):
    dest = tmp_path / 't.docx'
    officedoc.html_to_docx(
        '<table><tr><th>H1</th><th>H2</th></tr>'
        '<tr><td>a</td><td>b</td></tr></table>', str(dest))
    t = docx.Document(str(dest)).tables[0]
    assert t.cell(0, 0).paragraphs[0].runs[0].bold is True
    assert 'w:tblHeader' in t.rows[0]._tr.xml          # repeats on each page
    assert t.cell(1, 0).paragraphs[0].runs[0].bold is not True


def test_table_column_width_applied(tmp_path):
    from docx.shared import Inches
    dest = tmp_path / 't.docx'
    officedoc.html_to_docx(
        '<table><tr><td style="width:50%">half</td><td>rest</td></tr></table>',
        str(dest))
    t = docx.Document(str(dest)).tables[0]
    w = t.cell(0, 0).width
    assert w is not None and abs(w.inches - 3.25) < 0.05   # 50% of 6.5in


def test_table_cell_block_content_round_trips(tmp_path):
    """Cell text wrapped in <p> (editor / mammoth markup) must survive — a
    plain inline render would drop block children and leave cells empty."""
    dest = tmp_path / 't.docx'
    officedoc.html_to_docx(
        '<table><tr><th><p>Name</p></th><th><p>Qty</p></th></tr>'
        '<tr><td><p>Apple</p></td><td><p>3</p></td></tr></table>', str(dest))
    t = docx.Document(str(dest)).tables[0]
    assert t.cell(0, 0).text == 'Name'
    assert t.cell(1, 0).text == 'Apple' and t.cell(1, 1).text == '3'
    assert t.cell(0, 0).paragraphs[0].runs[0].bold is True


def test_table_cell_keeps_multiple_paragraphs(tmp_path):
    """A merged cell can hold several <p> blocks; each becomes its own line."""
    dest = tmp_path / 't.docx'
    officedoc.html_to_docx(
        '<table><tr><td><p>line one</p><p>line two</p></td></tr></table>', str(dest))
    cell = docx.Document(str(dest)).tables[0].cell(0, 0)
    assert [p.text for p in cell.paragraphs if p.text] == ['line one', 'line two']


# ── image editing round-trip (resize / align / alt text) ───────────────────────
_PNG = ('data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAIAAAACCAIAAAD91Jpz'
        'AAAAEElEQVR4nGP4z8AARAwQCgAf7gP9i18U1AAAAABJRU5ErkJggg==')


def test_image_width_alt_and_alignment_round_trip(tmp_path):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    dest = tmp_path / 'img.docx'
    officedoc.html_to_docx(
        f'<p style="text-align:center">'
        f'<img src="{_PNG}" style="width:50%" alt="A red dot"></p>', str(dest))
    d = docx.Document(str(dest))

    assert len(d.inline_shapes) == 1
    shape = d.inline_shapes[0]
    assert abs(shape.width.inches - 3.25) < 0.1            # 50% of 6.5in
    assert shape._inline.find(qn('wp:docPr')).get('descr') == 'A red dot'
    assert d.paragraphs[0].alignment == WD_ALIGN_PARAGRAPH.CENTER


def test_image_without_width_clamps_to_page(tmp_path):
    dest = tmp_path / 'img.docx'
    officedoc.html_to_docx(f'<p><img src="{_PNG}"></p>', str(dest))
    d = docx.Document(str(dest))
    assert len(d.inline_shapes) == 1                       # still embeds fine


# ── editor niceties (horizontal rule, Title/Subtitle paragraph styles) ─────────
def test_horizontal_rule_round_trips(tmp_path):
    """<hr> becomes an empty paragraph carrying a bottom border (a real rule)."""
    from docx.oxml.ns import qn
    dest = tmp_path / 'hr.docx'
    officedoc.html_to_docx('<p>above</p><hr><p>below</p>', str(dest))
    d = docx.Document(str(dest))
    borders = [p for p in d.paragraphs
               if p._p.find(qn('w:pPr')) is not None
               and p._p.find(qn('w:pPr')).find(qn('w:pBdr')) is not None]
    assert len(borders) == 1
    bottom = borders[0]._p.find(qn('w:pPr')).find(qn('w:pBdr')).find(qn('w:bottom'))
    assert bottom is not None and bottom.get(qn('w:val')) == 'single'


def test_title_and_subtitle_styles(tmp_path):
    """The editor's Title/Subtitle classes map to Word's Title/Subtitle styles."""
    dest = tmp_path / 'ts.docx'
    officedoc.html_to_docx(
        '<p class="doc-title">My Title</p>'
        '<p class="doc-subtitle">A subtitle</p>'
        '<p>body</p>', str(dest))
    d = docx.Document(str(dest))
    assert _by_text(d, 'My Title').style.name == 'Title'
    assert _by_text(d, 'A subtitle').style.name == 'Subtitle'
    assert _by_text(d, 'body').style.name == 'Normal'


def test_title_subtitle_open_side_maps_to_classes(tmp_path):
    """Opening a Word doc with Title/Subtitle styles surfaces editor classes,
    so the round-trip is symmetric (open → edit → save keeps the styles)."""
    src = docx.Document()
    src.add_paragraph('Doc Title', style='Title')
    src.add_paragraph('Doc Subtitle', style='Subtitle')
    src.add_paragraph('plain body')
    p = tmp_path / 'src.docx'
    src.save(str(p))
    html = officedoc._docx_to_html(str(p))['html']
    assert 'class="doc-title"' in html
    assert 'class="doc-subtitle"' in html


# ── page setup (size / orientation / margins) ──────────────────────────────────
def test_match_page_size_is_orientation_agnostic():
    assert officedoc._match_page_size(8.5, 11.0) == 'letter'
    assert officedoc._match_page_size(11.0, 8.5) == 'letter'   # rotated → same
    assert officedoc._match_page_size(8.27, 11.69) == 'a4'
    assert officedoc._match_page_size(8.5, 14.0) == 'legal'
    assert officedoc._match_page_size(5.0, 7.0) == 'custom'


def test_page_setup_round_trips(tmp_path):
    """Explicit dimensions + landscape + margins apply to the saved section."""
    from docx.enum.section import WD_ORIENT
    dest = tmp_path / 'pg.docx'
    page = {'size': 'a4', 'orientation': 'landscape',
            'width_in': 8.27, 'height_in': 11.69,
            'margins': {'top': 0.5, 'bottom': 0.5, 'left': 0.75, 'right': 0.75}}
    officedoc.html_to_docx('<p>hi</p>', str(dest), page=page)
    s = docx.Document(str(dest)).sections[0]
    assert s.orientation == WD_ORIENT.LANDSCAPE
    assert abs(s.page_width.inches - 11.69) < 0.02      # landscape swaps W/H
    assert abs(s.page_height.inches - 8.27) < 0.02
    assert abs(s.top_margin.inches - 0.5) < 0.01
    assert abs(s.left_margin.inches - 0.75) < 0.01


def test_page_setup_size_label_without_dims(tmp_path):
    """A bare size label (no explicit inches) still yields correct dimensions."""
    dest = tmp_path / 'lg.docx'
    officedoc.html_to_docx('<p>x</p>', str(dest),
                           page={'size': 'legal', 'orientation': 'portrait'})
    s = docx.Document(str(dest)).sections[0]
    assert abs(s.page_width.inches - 8.5) < 0.02
    assert abs(s.page_height.inches - 14.0) < 0.02


def test_html_to_docx_without_page_keeps_default_letter(tmp_path):
    """Omitting page leaves python-docx's default Letter-portrait geometry."""
    dest = tmp_path / 'd.docx'
    officedoc.html_to_docx('<p>hi</p>', str(dest))
    s = docx.Document(str(dest)).sections[0]
    assert abs(s.page_width.inches - 8.5) < 0.02
    assert abs(s.page_height.inches - 11.0) < 0.02


def test_page_setup_open_side_reads_section(tmp_path):
    """_docx_page_setup reports a section's size / orientation / margins."""
    from docx.shared import Inches
    from docx.enum.section import WD_ORIENT
    d = docx.Document()
    s = d.sections[0]
    s.orientation = WD_ORIENT.LANDSCAPE
    s.page_width = Inches(11.0)
    s.page_height = Inches(8.5)
    s.top_margin = Inches(0.5)
    p = tmp_path / 'pg.docx'
    d.save(str(p))
    ps = officedoc._docx_page_setup(str(p))
    assert ps is not None
    assert ps['orientation'] == 'landscape'
    assert ps['size'] == 'letter'                       # 8.5×11 normalised
    assert abs(ps['margins']['top'] - 0.5) < 0.01


# ── headers & footers (Tier 4b) ────────────────────────────────────────────────
def test_headers_footers_round_trip(tmp_path):
    """Header text + a footer page-number field survive the rebuild and read
    back symmetrically through the open-side reader."""
    dest = tmp_path / 'hf.docx'
    page = {'size': 'letter', 'orientation': 'portrait',
            'header': {'text': 'Confidential', 'page_num': False},
            'footer': {'text': 'Page', 'page_num': True}}
    officedoc.html_to_docx('<p>body</p>', str(dest), page=page)

    s = docx.Document(str(dest)).sections[0]
    assert s.header.is_linked_to_previous is False
    assert 'Confidential' in s.header.paragraphs[0].text
    fxml = s.footer.paragraphs[0]._p.xml
    assert 'PAGE' in fxml and 'fldSimple' in fxml       # a real page-number field

    hf = officedoc._docx_headers_footers(str(dest))
    assert hf['header']['text'] == 'Confidential'
    assert hf['header']['page_num'] is False
    assert hf['footer']['page_num'] is True


def test_header_multiline_round_trips(tmp_path):
    """A two-line header becomes two paragraphs in the header part."""
    dest = tmp_path / 'ml.docx'
    officedoc.html_to_docx('<p>body</p>', str(dest),
                           page={'header': {'text': 'Line one\nLine two'}})
    hdr = docx.Document(str(dest)).sections[0].header
    assert [p.text for p in hdr.paragraphs if p.text] == ['Line one', 'Line two']


def test_blank_headers_footers_leave_defaults(tmp_path):
    """Empty header/footer specs don't materialise content."""
    dest = tmp_path / 'blank.docx'
    officedoc.html_to_docx('<p>body</p>', str(dest),
                           page={'header': {'text': '', 'page_num': False},
                                 'footer': {'text': '', 'page_num': False}})
    s = docx.Document(str(dest)).sections[0]
    assert (s.header.paragraphs[0].text or '') == ''
    assert 'fldSimple' not in s.footer.paragraphs[0]._p.xml


# ── footnotes (Tier 4c) ──────────────────────────────────────────────────────
FN_BODY = (
    '<p>Intro<sup class="fn-ref" data-fn-id="1" contenteditable="false">1</sup>'
    ' and more<sup class="fn-ref" data-fn-id="2" contenteditable="false">2</sup>.</p>'
    '<section class="doc-footnotes" data-doc-footnotes="1"><ol class="fn-list">'
    '<li class="fn-item" data-fn-id="1">First note &amp; cite.</li>'
    '<li class="fn-item" data-fn-id="2">Second note.</li>'
    '</ol></section>')


def test_footnotes_round_trip_builds_a_real_part(tmp_path):
    """Editor footnote markup → a real footnotes part: content-type override, a
    document relationship, body reference marks and the (escaped) note text."""
    dest = tmp_path / 'fn.docx'
    officedoc.html_to_docx(FN_BODY, str(dest))
    with zipfile.ZipFile(str(dest)) as z:
        names = z.namelist()
        ct = z.read('[Content_Types].xml').decode('utf-8')
        rels = z.read('word/_rels/document.xml.rels').decode('utf-8')
        doc = z.read('word/document.xml').decode('utf-8')
        fn = z.read('word/footnotes.xml').decode('utf-8')
    assert 'word/footnotes.xml' in names
    assert 'footnotes+xml' in ct                       # content-type override
    assert 'relationships/footnotes' in rels           # document relationship
    assert doc.count('<w:footnoteReference') == 2       # two in-body marks
    assert 'First note &amp; cite.' in fn              # text present, & escaped
    assert 'Second note.' in fn


def test_footnotes_read_transform_and_fidelity(tmp_path):
    """A writer-produced doc reads back as inline markers + a footnotes section,
    and a plain-text footnote document is NOT flagged lossy."""
    dest = tmp_path / 'fn.docx'
    officedoc.html_to_docx(FN_BODY, str(dest))
    out = officedoc._docx_to_html(str(dest))
    assert 'class="fn-ref"' in out['html']
    assert 'doc-footnotes' in out['html']
    assert 'First note' in out['html']
    assert out['fidelity']['lossy'] is False
    assert 'Footnotes' not in out['fidelity']['features']


def test_footnote_refs_renumber_in_document_order(tmp_path):
    """Markers are renumbered 1..n by body order (not their original id), and the
    notes follow that order."""
    body = (
        '<p>A<sup class="fn-ref" data-fn-id="7">7</sup>'
        ' B<sup class="fn-ref" data-fn-id="3">3</sup></p>'
        '<section class="doc-footnotes"><ol class="fn-list">'
        '<li class="fn-item" data-fn-id="3">Note three.</li>'
        '<li class="fn-item" data-fn-id="7">Note seven.</li>'
        '</ol></section>')
    dest = tmp_path / 'fn.docx'
    officedoc.html_to_docx(body, str(dest))
    with zipfile.ZipFile(str(dest)) as z:
        doc = z.read('word/document.xml').decode('utf-8')
        fn = z.read('word/footnotes.xml').decode('utf-8')
    assert re.findall(r'<w:footnoteReference w:id="(\d+)"', doc) == ['1', '2']
    note1 = fn.split('w:id="1"')[1].split('</w:footnote>')[0]
    note2 = fn.split('w:id="2"')[1].split('</w:footnote>')[0]
    assert 'Note seven.' in note1                       # first-referenced note
    assert 'Note three.' in note2


def test_multiline_footnote_round_trips(tmp_path):
    """A footnote with two paragraphs becomes two w:p in the note."""
    body = (
        '<p>X<sup class="fn-ref" data-fn-id="1">1</sup></p>'
        '<section class="doc-footnotes"><ol class="fn-list">'
        '<li class="fn-item" data-fn-id="1"><p>Para one.</p><p>Para two.</p></li>'
        '</ol></section>')
    dest = tmp_path / 'fn.docx'
    officedoc.html_to_docx(body, str(dest))
    with zipfile.ZipFile(str(dest)) as z:
        fn = z.read('word/footnotes.xml').decode('utf-8')
    note = fn.split('w:id="1"')[1].split('</w:footnote>')[0]
    assert 'Para one.' in note and 'Para two.' in note
    assert note.count('<w:p>') == 2


def test_doc_without_footnotes_has_no_part(tmp_path):
    """Plain documents don't gain a footnotes part."""
    dest = tmp_path / 'plain.docx'
    officedoc.html_to_docx('<p>just text</p>', str(dest))
    with zipfile.ZipFile(str(dest)) as z:
        assert 'word/footnotes.xml' not in z.namelist()


def test_footnotes_are_simple_predicate():
    """Plain notes are simple; tables, drawings, hyperlinks or fields are not."""
    plain = ('<w:footnotes><w:footnote w:id="1"><w:p><w:r><w:t>hi</w:t>'
             '</w:r></w:p></w:footnote></w:footnotes>')
    assert officedoc._footnotes_are_simple(plain) is True
    assert officedoc._footnotes_are_simple(plain + '<w:tbl></w:tbl>') is False
    assert officedoc._footnotes_are_simple('<w:hyperlink/>') is False
    assert officedoc._footnotes_are_simple('<w:drawing/>') is False
    assert officedoc._footnotes_are_simple('<w:instrText> TIME </w:instrText>') is False
