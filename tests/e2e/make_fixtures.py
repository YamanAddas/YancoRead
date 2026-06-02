"""Generate the tiny sample documents the Playwright smoke tests open.

Writes into tests/e2e/fixtures/ (git-ignored). Run before `playwright test`
(the CI does this in a step; locally: `python tests/e2e/make_fixtures.py`).
"""
import os

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, 'fixtures')
os.makedirs(OUT, exist_ok=True)


def make_pdf(path):
    import fitz
    doc = fitz.open()
    for i in range(3):
        page = doc.new_page(width=612, height=792)
        page.insert_text((72, 96), f'Smoke test page {i + 1} — hello world.', fontsize=18)
    doc.save(path)
    doc.close()


def make_docx(path):
    from docx import Document
    d = Document()
    d.add_heading('Smoke Heading', level=1)
    d.add_paragraph('A paragraph the office reader should render.')
    d.save(path)


def make_xlsx(path):
    from openpyxl import Workbook
    wb = Workbook(); ws = wb.active; ws.title = 'Sheet1'
    ws.append(['Region', 'Value'])
    ws.append(['EU', 1200])
    wb.save(path)


def make_pptx(path):
    from pptx import Presentation
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = 'Smoke Slide'
    s.placeholders[1].text = 'A slide body.'
    prs.save(path)


def make_txt(path):
    with open(path, 'w', encoding='utf-8') as f:
        f.write('Plain text smoke fixture.\n' + 'Line of content.\n' * 5)


def make_png(path):
    from PIL import Image
    Image.new('RGB', (320, 240), (240, 244, 250)).save(path, 'PNG')


if __name__ == '__main__':
    make_pdf(os.path.join(OUT, 'sample.pdf'))
    make_docx(os.path.join(OUT, 'sample.docx'))
    make_xlsx(os.path.join(OUT, 'sample.xlsx'))
    make_pptx(os.path.join(OUT, 'sample.pptx'))
    make_txt(os.path.join(OUT, 'sample.txt'))
    make_png(os.path.join(OUT, 'sample.png'))
    print('fixtures written to', OUT)
